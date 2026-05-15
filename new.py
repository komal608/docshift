import os
import io
from queue import Full
import uuid
import json
import time
import glob
import base64
import tempfile
import logging
import requests
import sqlite3
import zipfile
import subprocess
import re
import mimetypes
import random
import threading
import smtplib
import csv
import xlsxwriter
import html
import fitz  # PyMuPDF
import firebase_admin
from firebase_admin import credentials, db, auth
import cloudinary
import cloudinary.uploader
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from datetime import datetime, timedelta
from functools import wraps
from collections import deque
from dotenv import load_dotenv
from flask import Flask, request, jsonify, render_template, session, redirect, url_for, send_file, send_from_directory
import difflib
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.shared import Pt
from flask_cors import CORS
from PIL import Image
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from docx import Document
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from gtts import gTTS
import speech_recognition as sr
from pydub import AudioSegment
from bs4 import BeautifulSoup
from fpdf import FPDF


# --- Flask & Environment Setup ---

app = Flask(__name__)
load_dotenv('api.env')

# Firebase Realtime Database config
FIREBASE_CRED_PATH = os.getenv('FIREBASE_CRED_PATH')
FIREBASE_DB_URL = 'https://docshift-86065-default-rtdb.firebaseio.com/'

if not firebase_admin._apps:
    cred = credentials.Certificate(json.loads(FIREBASE_CRED_PATH))
    firebase_admin.initialize_app(cred, {'databaseURL': FIREBASE_DB_URL})

# Cloudinary config
cloudinary.config(
    cloud_name='dvdeflyta',
    api_key='568435982421747',
    api_secret='-xqsm00d0D9Hxp1YsrA6OrU-hpw'
)

# Flask CORS and secret key
CORS(app)
app.secret_key = os.urandom(24)

# Logging config
logging.basicConfig(level=logging.DEBUG,
                    format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# OpenRouter API settings (change with your key)
OPENROUTER_API_KEY = os.getenv('OPENROUTER_API_KEY', 'your_api_key_here')
OPENROUTER_API_URL = 'https://openrouter.ai/api/v1/chat/completions'
OPENROUTER_MODEL = 'gpt-4o-mini'

# Speech recognizer init
recognizer = sr.Recognizer()

# Globals for AI Document Screener and AI PDF Editor
current_document_text = ''
# Remove chat history limit for resume analyzer chat (unlimited turns)
conversation_history = []
latest_text = ""

# --- Utility Functions ---

def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'logged_in' not in session:
            # Check if this is an AJAX request
            if request.headers.get('Content-Type') == 'application/json' or request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({'error': 'Authentication required'}), 401
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

# --- AI Resume Analyzer Utilities ---
def extract_text_from_file(file_storage):
    import pytesseract
    from pdf2image import convert_from_bytes
    """Extract text from uploaded file (PDF, DOCX, TXT, etc.) with OCR fallback for image-based PDFs."""
    filename = file_storage.filename.lower()
    ext = os.path.splitext(filename)[1]
    file_storage.seek(0)
    
    if ext == '.pdf':
        try:
            # Step 1: Try PyPDF2 for text-based PDF
            file_storage.seek(0)
            reader = PdfReader(file_storage)
            extracted_text = "\n".join(page.extract_text() or '' for page in reader.pages)
            extracted_text = extracted_text.strip()
            
            # Check if extraction is minimal (image-based PDF)
            text_length = len(extracted_text.replace(" ", "").replace("\n", ""))
            
            # Step 2: If minimal text, use OCR as fallback
            if text_length < 100:
                logger.info("Minimal text from PyPDF2. Attempting OCR extraction...")
                try:
                    file_storage.seek(0)
                    pdf_bytes = file_storage.read()
                    images = convert_from_bytes(pdf_bytes, first_page=1, last_page=min(5, len(reader.pages)))
                    
                    ocr_text = ""
                    for page_num, image in enumerate(images):
                        page_ocr = pytesseract.image_to_string(image)
                        ocr_text += f"\n{page_ocr}\n"
                    
                    extracted_text = ocr_text.strip()
                    logger.info(f"OCR extraction successful: {len(extracted_text)} characters")
                except Exception as ocr_e:
                    logger.warning(f"OCR extraction failed: {ocr_e}. Using PyPDF2 result.")
            
            return extracted_text if extracted_text else "[No text found in PDF]"
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            return f"[PDF extraction error: {e}]"
    
    elif ext in ['.docx']:
        try:
            doc = Document(file_storage)
            return "\n".join([p.text for p in doc.paragraphs]).strip()
        except Exception as e:
            return f"[DOCX extraction error: {e}]"
    
    elif ext in ['.txt']:
        try:
            return file_storage.read().decode(errors='ignore')
        except Exception as e:
            return f"[TXT extraction error: {e}]"
    
    else:
        return "[Unsupported file type]"

def summarize_resume_with_openrouter(text):
    """Send resume text to OpenRouter API and get a 200-word summary."""
    prompt = f"You are a highly precise AI Resume Analyzer. Your task is to deeply analyze the entire resume provided to you by carefully reading every line, section, heading, and detail without skipping anything. The resume is the single source of truth. You must extract and understand all information including name, contact details, college or university, degree, specialization, HSC and SSC details if present, CGPA or percentage, skills, internships, work experience, projects, certifications, achievements, tools, technologies, and any other information mentioned. When answering questions, respond strictly and completely based only on the resume content. If a detail exists anywhere in the resume, you must provide it clearly and accurately. Never guess, assume, or fabricate information. If and only if the requested information is truly not present in the resume, reply exactly with: 'This information is not mentioned in the resume.' Your goal is to ensure that even the smallest detail in the resume is correctly analyzed and answered.\n\n{text}"
    headers = {
        'Authorization': f'Bearer {OPENROUTER_API_KEY}',
        'Content-Type': 'application/json',
    }
    data = {
        "model": OPENROUTER_MODEL,
        "messages": [
            {"role": "system", "content": "You are a professional resume summarizer."},
            {"role": "user", "content": prompt}
        ],
        "max_tokens": 800,
        "temperature": 0.5
    }
    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data, timeout=60)
        response.raise_for_status()
        result = response.json()
        summary = result.get('choices', [{}])[0].get('message', {}).get('content', '').strip()
        return summary or '[No summary returned]'
    except Exception as e:
        return f"[OpenRouter API error: {e}]"

# --- Flask Route: AI Resume Analyzer ---
@app.route('/ai_resume_analyzer', methods=['POST'])
@login_required
def ai_resume_analyzer():
    """Upload multiple resumes, analyze, and return 200-word summaries."""
    if 'resumes' not in request.files:
        return jsonify({"error": "No files part in the request."}), 400
    files = request.files.getlist('resumes')
    results = []
    for file_storage in files:
        filename = secure_filename(file_storage.filename)
        text = extract_text_from_file(file_storage)
        summary = summarize_resume_with_openrouter(text)
        results.append({
            "filename": filename,
            "summary": summary
        })
    # Store results in session for download
    session['resume_analyzer_results'] = results
    return jsonify({"results": results})

# Render the AI Resume Analyzer page
@app.route('/ai_resume_analyzer', methods=['GET'])
@login_required
def ai_resume_analyzer_page():
    # Pass user context for navbar/profile dropdown
    return render_template('ai_resume_analyzer.html', **get_user_context())

# --- Flask Route: AI Resume Analyzer Chat ---
@app.route('/ai_resume_analyzer_chat', methods=['POST'])
@login_required
def ai_resume_analyzer_chat():
    """Chat about the analyzed resumes (summaries in session). Unlimited chat history."""
    data = request.get_json()
    message = data.get('message', '')
    history = data.get('history', [])
    summaries = session.get('resume_analyzer_results', [])
    if not summaries:
        return jsonify({'reply': 'No resumes have been analyzed yet.'})
    # Compose context for chat: all summaries
    context = '\n\n'.join(f"Resume: {item['filename']}\nSummary: {item['summary']}" for item in summaries)
    # Build chat history for OpenRouter (unlimited turns)
    chat_messages = [
        {"role": "system", "content": "You are an expert HR assistant. Answer questions based only on the following analyzed resume summaries. If the answer is not in the summaries, say you don't know. You can compare, rank, and analyze the resumes as requested.\n\n" + context}
    ]
    # Append all previous chat turns (no limit)
    for msg in history:
        if msg['role'] == 'user':
            chat_messages.append({"role": "user", "content": msg['content']})
        elif msg['role'] == 'assistant':
            chat_messages.append({"role": "assistant", "content": msg['content']})
    chat_messages.append({"role": "user", "content": message})
    headers = {
        'Authorization': f'Bearer {OPENROUTER_API_KEY}',
        'Content-Type': 'application/json',
    }
    data = {
        "model": OPENROUTER_MODEL,
        "messages": chat_messages,
        "max_tokens": 1200,  # Increase token limit for more complex answers
        "temperature": 0.5
    }
    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data, timeout=90)
        response.raise_for_status()
        result = response.json()
        reply = result.get('choices', [{}])[0].get('message', {}).get('content', '').strip()
        return jsonify({'reply': reply or '[No response from AI]'}), 200
    except Exception as e:
        return jsonify({'reply': f'[OpenRouter API error: {e}]'}), 500

# Download summaries as PDF or Excel, then upload to Firebase/Cloudinary
@app.route('/download_resume_summaries')
@login_required
def download_resume_summaries():
    import io
    from fpdf import FPDF
    import xlsxwriter
    results = session.get('resume_analyzer_results', [])
    fmt = request.args.get('format', 'pdf')
    username = session.get('username', 'admin')
    output = io.BytesIO()
    filename = f"resume_summaries_{int(time.time())}.{fmt if fmt=='pdf' else 'xlsx'}"
    if fmt == 'pdf':
        def to_latin1(text):
            return text.encode('latin1', errors='replace').decode('latin1')
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font('Arial', 'B', 16)
        pdf.cell(0, 10, to_latin1('AI Resume Analyzer Summaries'), ln=True, align='C')
        pdf.ln(10)
        pdf.set_font('Arial', '', 12)
        for res in results:
            pdf.set_font('Arial', 'B', 12)
            pdf.cell(0, 10, to_latin1(res['filename']), ln=True)
            pdf.set_font('Arial', '', 11)
            pdf.multi_cell(0, 8, to_latin1(res['summary']))
            pdf.ln(5)
        pdf_bytes = pdf.output(dest='S').encode('latin1', errors='replace')
        output.write(pdf_bytes)
        output.seek(0)
        mimetype = 'application/pdf'
    else:
        workbook = xlsxwriter.Workbook(output, {'in_memory': True})
        worksheet = workbook.add_worksheet('Summaries')
        worksheet.write(0, 0, 'Filename')
        worksheet.write(0, 1, 'Summary')
        for idx, res in enumerate(results, 1):
            worksheet.write(idx, 0, res['filename'])
            worksheet.write(idx, 1, res['summary'])
        workbook.close()
        output.seek(0)
        mimetype = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    # Upload to Cloudinary
    output.seek(0)
    cloudinary_result = cloudinary.uploader.upload(output, resource_type='raw', folder='resume_analyzer', public_id=filename)
    url = cloudinary_result.get('secure_url')
    # Store URL in Firebase under user's storage
    store_url_in_firebase(url, 'resume_analyzer', filename)
    # Send file for download
    output.seek(0)
    return send_file(output, as_attachment=True, download_name=filename, mimetype=mimetype)

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'logged_in' not in session:
            return redirect(url_for('login'))
        if session.get('role') != 'admin':
            return redirect(url_for('index'))
        return f(*args, **kwargs)
    return decorated_function

def store_url_in_firebase(url, category, filename):
    """Store URL in firebase per user under storage/{username}/{category}."""
    safe_key = re.sub(r'[./#$\[\]]', '_', filename)
    username = session.get('username', 'admin')
    ref = db.reference(f'storage/{username}/{category}/{safe_key}')
    ref.set({'filename': filename, 'url': url})
    return True

# --- Phone & Email Verification Functions ---
"""
--- AI Notes Generator ---
"""

def extract_text_for_notes(file_storage):
    """Extract raw text from PDF, DOCX, or TXT for notes generation."""
    filename = file_storage.filename.lower()
    ext = os.path.splitext(filename)[1]
    file_storage.seek(0)
    if ext == '.pdf':
        try:
            import fitz
            doc = fitz.open(stream=file_storage.read(), filetype="pdf")
            text = "".join(page.get_text() for page in doc)
            return text.strip()
        except Exception as e:
            return f"[PDF extraction error: {e}]"
    elif ext == '.docx':
        try:
            from docx import Document
            file_storage.seek(0)
            doc = Document(file_storage)
            return "\n".join([p.text for p in doc.paragraphs]).strip()
        except Exception as e:
            return f"[DOCX extraction error: {e}]"
    elif ext == '.txt':
        try:
            return file_storage.read().decode(errors='ignore')
        except Exception as e:
            return f"[TXT extraction error: {e}]"
    else:
        return "[Unsupported file type]"

def save_notes_to_docx(notes: str, output_path: str):
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    for line in notes.split("\n"):
        if line.startswith("## "):
            doc.add_heading(line.replace("## ", "").strip(), level=1)
        elif line.startswith("- "):
            doc.add_paragraph(line.replace("- ", "").strip(), style="List Bullet")
        else:
            text = line.strip()
            if text:
                p = doc.add_paragraph(text)
                if p.runs:
                    p.runs[0].font.size = Pt(11)
    doc.save(output_path)

def save_notes_to_pdf(notes: str, output_path: str):
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    x, y = 50, height - 50
    for line in notes.split("\n"):
        if line.startswith("## "):
            c.setFont("Helvetica-Bold", 14)
            c.drawString(x, y, line.replace("## ", "").strip())
        elif line.startswith("- "):
            c.setFont("Helvetica", 11)
            c.drawString(x + 20, y, "• " + line.replace("- ", "").strip())
        else:
            c.setFont("Helvetica", 11)
            c.drawString(x, y, line.strip())
        y -= 18
        if y < 50:
            c.showPage()
            y = height - 50
    c.save()

def generate_notes_with_openrouter(raw_text):
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": OPENROUTER_MODEL,
        "messages": [
            {"role": "system", "content": "You are an assistant that creates concise, structured, point-wise notes for documents. Use bullet points, numbers, or clear sections, but do NOT use any hashtags (#), markdown headings, or heading symbols. Only output clean, readable notes."},
            {"role": "user", "content": f"Document content:\n\n{raw_text}\n\nSummarize this into structured, point-wise notes with bullets or numbers, but do not use any hashtags or markdown heading symbols anywhere in your answer."}
        ]
    }
    response = requests.post(OPENROUTER_API_URL, headers=headers, json=payload)
    if response.status_code == 200:
        result = response.json()
        notes = result["choices"][0]["message"]["content"]
        return notes.strip()
    else:
        raise Exception(f"Error: {response.status_code} - {response.text}")

# --- Flask Route: AI Notes Generator ---
@app.route('/ai_notes_generator', methods=['POST'])
@login_required
def ai_notes_generator():
    """Upload a document, generate AI notes, and return download links."""
    if 'document' not in request.files:
        return jsonify({"error": "No file part in the request."}), 400
    file_storage = request.files['document']
    filename = secure_filename(file_storage.filename)
    raw_text = extract_text_for_notes(file_storage)
    try:
        notes = generate_notes_with_openrouter(raw_text)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    # Save notes to files (no upload yet)
    output_dir = os.path.join(tempfile.gettempdir(), f"notes_{uuid.uuid4().hex}")
    os.makedirs(output_dir, exist_ok=True)
    docx_path = os.path.join(output_dir, "notes.docx")
    pdf_path = os.path.join(output_dir, "notes.pdf")
    save_notes_to_docx(notes, docx_path)
    save_notes_to_pdf(notes, pdf_path)

    # Store results in session for download (no upload yet)
    session['ai_notes_generator_results'] = {
        "notes": notes,
        "docx_file": docx_path,
        "pdf_file": pdf_path,
        "filename": filename
    }
    return jsonify({"notes": notes, "filename": filename})

# --- Flask Route: Download AI Notes PDF ---
@app.route('/download_ai_notes_pdf', methods=['GET'])
@login_required
def download_ai_notes_pdf():
    results = session.get('ai_notes_generator_results')
    if not results or not os.path.exists(results['pdf_file']):
        return jsonify({"error": "No PDF file found. Please generate notes first."}), 400
    filename = results['filename']
    pdf_path = results['pdf_file']
    # Upload to Cloudinary
    with open(pdf_path, "rb") as pdf_file:
        pdf_upload = cloudinary.uploader.upload(pdf_file, resource_type='raw', folder='ai_notes_generator', public_id=f"{filename}_notes_pdf")
    pdf_url = pdf_upload.get('secure_url')
    # Store URL in Firebase
    store_url_in_firebase(pdf_url, 'ai_notes_generator', f"{filename}_notes.pdf")
    # Send file for download
    from flask import send_file
    return send_file(pdf_path, as_attachment=True, download_name=f"{filename}_notes.pdf", mimetype='application/pdf')

# --- Flask Route: Download AI Notes DOCX ---
@app.route('/download_ai_notes_docx', methods=['GET'])
@login_required
def download_ai_notes_docx():
    results = session.get('ai_notes_generator_results')
    if not results or not os.path.exists(results['docx_file']):
        return jsonify({"error": "No DOCX file found. Please generate notes first."}), 400
    filename = results['filename']
    docx_path = results['docx_file']
    # Upload to Cloudinary
    with open(docx_path, "rb") as docx_file:
        docx_upload = cloudinary.uploader.upload(docx_file, resource_type='raw', folder='ai_notes_generator', public_id=f"{filename}_notes_docx")
    docx_url = docx_upload.get('secure_url')
    # Store URL in Firebase
    store_url_in_firebase(docx_url, 'ai_notes_generator', f"{filename}_notes.docx")
    # Send file for download
    from flask import send_file
    return send_file(docx_path, as_attachment=True, download_name=f"{filename}_notes.docx", mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
# Render the AI Notes Generator page
@app.route('/ai_notes_generator', methods=['GET'])
@login_required
def ai_notes_generator_page():
    return render_template('ai_notes_generator.html', **get_user_context())

def is_production_mode():
    """Check if we're running in production mode with real credentials"""
    smtp_configured = bool(os.getenv('SMTP_USERNAME') and os.getenv('SMTP_PASSWORD'))
    return smtp_configured

def generate_otp():
    """Generate a random 6-digit OTP"""
    return str(random.randint(100000, 999999))

def store_otp_in_firebase(identifier, otp_code, username, verification_type='phone'):
    """Store OTP in Firebase with expiration (10 minutes)"""
    try:
        # Create expiration timestamp (10 minutes from now)
        expiry_time = datetime.now() + timedelta(minutes=10)
        expiry_timestamp = int(expiry_time.timestamp())
        
        # Store OTP data
        otp_data = {
            'code': otp_code,
            'type': verification_type,
            'expires_at': expiry_timestamp,
            'username': username,
            'verified': False,
            'created_at': int(datetime.now().timestamp())
        }
        
        # Clean identifier for Firebase key
        if verification_type == 'phone':
            clean_key = re.sub(r'[^0-9]', '', identifier)
        else:  # email
            clean_key = re.sub(r'[./#$\[\]@]', '_', identifier)
        
        ref_path = f'verification_codes/{verification_type}_{clean_key}'
        logger.info(f"Storing OTP at path: {ref_path} | Original {verification_type}: {identifier} | Clean key: {clean_key}")
        
        ref = db.reference(ref_path)
        ref.set(otp_data)
        
        logger.info(f"OTP stored successfully for {verification_type} {identifier}: {otp_code}")
        return True
    except Exception as e:
        logger.error(f"Error storing OTP: {str(e)}")
        return False

def store_phone_otp_in_firebase(phone_number, otp_code, username):
    """Store phone OTP in Firebase - backward compatibility"""
    return store_otp_in_firebase(phone_number, otp_code, username, 'phone')

def store_email_otp_in_firebase(email, otp_code, username):
    """Store email OTP in Firebase"""
    return store_otp_in_firebase(email, otp_code, username, 'email')

def verify_otp_from_firebase(identifier, submitted_otp, verification_type='phone'):
    """Verify OTP from Firebase"""
    try:
        if verification_type == 'phone':
            clean_key = re.sub(r'[^0-9]', '', identifier)
        else:  # email
            clean_key = re.sub(r'[./#$\[\]@]', '_', identifier)
        
        ref_path = f'verification_codes/{verification_type}_{clean_key}'
        logger.info(f"Looking for OTP at path: {ref_path} | Original {verification_type}: {identifier} | Clean key: {clean_key}")
        
        ref = db.reference(ref_path)
        stored_data = ref.get()
        
        if not stored_data:
            logger.warning(f"No OTP found for {verification_type}: {identifier} at path {ref_path}")
            return False, f"No OTP found for this {verification_type}. Please request a new code."
        
        # Check if OTP has expired
        current_timestamp = int(datetime.now().timestamp())
        if current_timestamp > stored_data.get('expires_at', 0):
            # Clean up expired OTP
            logger.warning(f"OTP expired for {verification_type}: {identifier}")
            ref.delete()
            return False, "OTP has expired. Please request a new code."
        
        # Check if OTP matches
        if stored_data.get('code') != submitted_otp:
            logger.warning(f"Invalid OTP submitted for {verification_type}: {identifier}")
            return False, "Invalid OTP. Please check and try again."
        
        # Mark as verified and clean up
        stored_data['verified'] = True
        stored_data['verified_at'] = current_timestamp
        ref.set(stored_data)
        
        # Update user profile to mark as verified
        username = stored_data.get('username')
        if username:
            user_ref = db.reference(f'Data/{username}')
            user_data = user_ref.get() or {}
            
            if verification_type == 'phone':
                user_data.update({
                    'phone_verified': True,
                    'phone_verified_at': datetime.now().isoformat()
                })
            else:  # email
                user_data.update({
                    'email_verified': True,
                    'email_verified_at': datetime.now().isoformat()
                })
            user_ref.set(user_data)
            logger.info(f"Marked {verification_type} as verified for username: {username}")
        
        # Clean up OTP after successful verification
        threading.Timer(5.0, lambda: ref.delete()).start()
        
        logger.info(f"{verification_type.title()} verified successfully for {identifier}")
        return True, f"{verification_type.title()} verified successfully!"
        
    except Exception as e:
        logger.error(f"Error verifying OTP for {verification_type} {identifier}: {str(e)}")
        return False, f"Verification error: {str(e)}"

def verify_phone_otp_from_firebase(phone_number, submitted_otp):
    """Verify phone OTP - backward compatibility"""
    return verify_otp_from_firebase(phone_number, submitted_otp, 'phone')

def verify_email_otp_from_firebase(email, submitted_otp):
    """Verify email OTP"""
    return verify_otp_from_firebase(email, submitted_otp, 'email')

def send_email_otp(email, otp_code):
    """Send email OTP using Gmail SMTP"""
    try:
        # Get email credentials from environment
        smtp_server = os.getenv('SMTP_SERVER', 'smtp.gmail.com')
        smtp_port = int(os.getenv('SMTP_PORT', '587'))
        smtp_username = os.getenv('SMTP_USERNAME')  # Your Gmail address
        smtp_password = os.getenv('SMTP_PASSWORD')  # Your Gmail app password
        
        # Check if email credentials are configured
        if not smtp_username or not smtp_password:
            logger.error("📧 Email credentials not configured")
            return False, "Email service not configured. Please contact administrator."
        
        # Create email message
        msg = MIMEMultipart()
        msg['From'] = smtp_username
        msg['To'] = email
        msg['Subject'] = "DocShift - Email Verification Code"
        
        # Create HTML email body
        html_body = f"""
        <html>
            <body style="font-family: Arial, sans-serif; line-height: 1.6; color: #333;">
                <div style="max-width: 600px; margin: 0 auto; padding: 20px; border: 1px solid #ddd; border-radius: 10px;">
                    <div style="text-align: center; margin-bottom: 30px;">
                        <h1 style="color: #4A90E2;">DocShift</h1>
                        <h2 style="color: #333;">Email Verification</h2>
                    </div>
                    
                    <div style="background-color: #f9f9f9; padding: 20px; border-radius: 5px; margin-bottom: 20px;">
                        <p>Hello,</p>
                        <p>You've requested to verify your email address for DocShift. Please use the following verification code:</p>
                        
                        <div style="text-align: center; margin: 30px 0;">
                            <span style="background-color: #4A90E2; color: white; padding: 15px 30px; font-size: 24px; font-weight: bold; border-radius: 5px; letter-spacing: 3px;">{otp_code}</span>
                        </div>
                        
                        <p><strong>This code will expire in 10 minutes.</strong></p>
                        <p>If you didn't request this verification, please ignore this email.</p>
                    </div>
                    
                    <div style="text-align: center; color: #666; font-size: 12px;">
                        <p>© 2025 DocShift. All rights reserved.</p>
                    </div>
                </div>
            </body>
        </html>
        """
        
        msg.attach(MIMEText(html_body, 'html'))
        
        # Send email
        server = smtplib.SMTP(smtp_server, smtp_port)
        server.starttls()
        server.login(smtp_username, smtp_password)
        text = msg.as_string()
        server.sendmail(smtp_username, email, text)
        server.quit()
        
        logger.info(f"📧 Email OTP sent successfully to {email}")
        return True, "Verification code sent to your email!"
        
    except Exception as e:
        logger.error(f"Error sending email: {str(e)}")
        return False, f"Failed to send email verification code. Please try again."

# Phone verification via email functionality removed

# SMS OTP functionality removed

def cleanup_expired_otps():
    """Clean up expired OTPs from Firebase"""
    try:
        ref = db.reference('verification_codes')
        all_codes = ref.get() or {}
        current_timestamp = int(datetime.now().timestamp())
        
        for key, data in all_codes.items():
            if isinstance(data, dict) and current_timestamp > data.get('expires_at', 0):
                db.reference(f'verification_codes/{key}').delete()
                logger.info(f"Cleaned up expired OTP for {key}")
                
    except Exception as e:
        logger.error(f"Error cleaning up expired OTPs: {str(e)}")

# Start background cleanup thread for expired OTPs
def periodic_otp_cleanup():
    while True:
        time.sleep(300)  # Check every 5 minutes
        cleanup_expired_otps()

cleanup_thread = threading.Thread(target=periodic_otp_cleanup, daemon=True)
cleanup_thread.start()

# --- End Phone & Email Verification Functions ---

def upload_to_cloudinary(local_path, folder):
    """Upload file to Cloudinary, handle image/raw types."""
    ext = os.path.splitext(local_path)[1].lower()
    if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp']:
        response = cloudinary.uploader.upload(local_path, folder=folder)
    elif ext in ['.mp3', '.wav', '.aac', '.ogg', '.flac', '.pdf', '.txt', '.docx', '.xlsx', '.pptx', '.csv']:
        response = cloudinary.uploader.upload(local_path, folder=folder, resource_type='raw')
    else:
        response = cloudinary.uploader.upload(local_path, folder=folder, resource_type='raw')
    return response['secure_url']

def get_user_storage_path():
    return f"storage/{session['username']}"

def init_db():
    """Initialize SQLite DB for conversion logs (optional)."""
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS conversions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        conversion_type TEXT NOT NULL,
        original_filename TEXT NOT NULL,
        converted_filename TEXT NOT NULL,
        timestamp DATETIME DEFAULT CURRENT_TIMESTAMP,
        file_path TEXT,
        cloudinary_url TEXT,
        username TEXT,
        status TEXT DEFAULT 'error'
    )''')
    # Add new columns if they don't exist (for existing databases)
    try:
        c.execute('ALTER TABLE conversions ADD COLUMN cloudinary_url TEXT')
    except sqlite3.OperationalError:
        pass  # Column already exists
    try:
        c.execute('ALTER TABLE conversions ADD COLUMN username TEXT')
    except sqlite3.OperationalError:
        pass  # Column already exists
    try:
        c.execute('ALTER TABLE conversions ADD COLUMN status TEXT DEFAULT "error"')
    except sqlite3.OperationalError:
        pass  # Column already exists
    c.execute('''CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT NOT NULL UNIQUE,
        password_hash TEXT NOT NULL
    )''')
    conn.commit()
    conn.close()

def log_conversion(conversion_type, original_filename, converted_filename, file_path=None, cloudinary_url=None, status=None):
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    username = session.get('username', 'admin')
    
    # Determine status based on cloudinary_url if not explicitly provided
    if status is None:
        status = 'success' if cloudinary_url else 'error'
    
    c.execute('''
        INSERT INTO conversions (conversion_type, original_filename, converted_filename, file_path, cloudinary_url, username, status)
        VALUES (?, ?, ?, ?, ?, ?, ?)
    ''', (conversion_type, original_filename, converted_filename, file_path, cloudinary_url, username, status))
    conn.commit()
    conn.close()

def is_ghostscript_installed():
    try:
        subprocess.run(['gs', '--version'], stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
        return True
    except (subprocess.CalledProcessError, FileNotFoundError):
        return False

# --- Firebase: Get user credentials ---
def get_user_by_username(username):
    """Fetch user credentials from Firebase Realtime DB.
    Assign role 'admin' for admin user, 'user' for regular users."""
    if username == 'admin':
        ref = db.reference('credentials/admin/admin')
        user_record = ref.get()
        if user_record:
            user_record['role'] = 'admin'
        return user_record
    else:
        ref = db.reference(f'credentials/users/{username}')
        user_record = ref.get()
        if user_record:
            user_record['role'] = 'user'
        return user_record

# --- Ensure admin credentials exist ---
def ensure_admin_credentials():
    admin_ref = db.reference('credentials/admin/admin')
    if not admin_ref.get():
        default_password = generate_password_hash("admin123")
        admin_ref.set({'password': default_password})
        print("Default admin credentials inserted (admin/admin123)")
ensure_admin_credentials()

# --- Routes ---

# --- Phone & Email Verification Routes ---

@app.route('/send_phone_otp', methods=['POST'])
def send_phone_otp():
    """Phone verification disabled - use email verification only"""
    return jsonify({'success': False, 'error': 'Phone verification is currently disabled. Please use email verification.'})

@app.route('/send_email_otp', methods=['POST'])
def send_email_otp_route():
    """Send OTP to email for verification"""
    try:
        data = request.get_json()
        email = data.get('email', '').strip()
        username = session.get('username')
        
        if not email:
            return jsonify({'success': False, 'error': 'Email address is required'})
        
        if not username:
            return jsonify({'success': False, 'error': 'User not logged in'})
        
        # Validate email format (basic validation)
        import re
        email_pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
        if not re.match(email_pattern, email):
            return jsonify({'success': False, 'error': 'Please enter a valid email address'})
        
        # Generate OTP
        otp_code = generate_otp()
        
        # Store OTP in Firebase
        if not store_email_otp_in_firebase(email, otp_code, username):
            return jsonify({'success': False, 'error': 'Failed to generate OTP'})
        
        # Send Email - production mode only
        email_success, email_message = send_email_otp(email, otp_code)
        
        if email_success:
            return jsonify({
                'success': True, 
                'message': 'Verification code sent to your email address. Please check your inbox.'
            })
        else:
            return jsonify({'success': False, 'error': email_message})
            
    except Exception as e:
        logger.error(f"Send email OTP error: {str(e)}")
        return jsonify({'success': False, 'error': 'Failed to send OTP'})

@app.route('/verify_phone_otp', methods=['POST'])
def verify_phone_otp():
    """Phone verification disabled - use email verification only"""
    return jsonify({'success': False, 'error': 'Phone verification is currently disabled. Please use email verification.'})

@app.route('/verify_email_otp', methods=['POST'])
def verify_email_otp():
    """Verify OTP for email"""
    try:
        data = request.get_json()
        email = data.get('email', '').strip()
        otp_code = data.get('otp_code', '').strip()
        
        if not email or not otp_code:
            return jsonify({'success': False, 'error': 'Email and OTP are required'})
        
        # Verify OTP
        verification_success, message = verify_email_otp_from_firebase(email, otp_code)
        
        if verification_success:
            return jsonify({
                'success': True,
                'message': message
            })
        else:
            return jsonify({
                'success': False,
                'error': message
            })
            
    except Exception as e:
        logger.error(f"Verify email OTP error: {str(e)}")
        return jsonify({'success': False, 'error': 'Failed to verify OTP'})

@app.route('/check_phone_verification_status', methods=['POST'])
def check_phone_verification_status():
    """Phone verification disabled - always return unverified"""
    return jsonify({'verified': False, 'message': 'Phone verification is currently disabled'})

@app.route('/check_email_verification_status', methods=['POST'])
def check_email_verification_status():
    """Check if email is verified for current user"""
    try:
        username = session.get('username')
        if not username:
            return jsonify({'verified': False, 'error': 'User not logged in'})
        
        user_ref = db.reference(f'Data/{username}')
        user_data = user_ref.get() or {}
        
        is_verified = user_data.get('email_verified', False)
        verified_at = user_data.get('email_verified_at', None)
        
        return jsonify({
            'verified': is_verified,
            'verified_at': verified_at
        })
        
    except Exception as e:
        logger.error(f"Check email verification status error: {str(e)}")
        return jsonify({'verified': False, 'error': 'Failed to check verification status'})

@app.route('/verify_email_standalone', methods=['POST'])
def verify_email_standalone():
    """Standalone email verification route (if needed for separate email verification page)"""
    try:
        email = request.form.get('email', '').strip()
        email_otp = request.form.get('email_otp', '').strip()
        
        if not email or not email_otp:
            return render_template('verify_email.html', 
                                 email=email,
                                 error='Email and OTP are required')
        
        # Verify OTP
        verification_success, message = verify_email_otp_from_firebase(email, email_otp)
        
        if verification_success:
            return render_template('registration_success.html', 
                                 message='Email verified successfully!')
        else:
            return render_template('verify_email.html', 
                                 email=email,
                                 error=message)
    
    except Exception as e:
        logger.error(f"Standalone email verification error: {str(e)}")
        return render_template('verify_email.html', 
                             email=email,
                             error='An error occurred during verification. Please try again.')

# --- End Phone & Email Verification Routes ---

@app.route('/register-company', methods=['GET', 'POST'])
def register_company():
    if request.method == 'POST':
        company_name = request.form.get('company_name')
        owner_name = request.form.get('owner_name')
        email = request.form.get('email')
        phone = request.form.get('phone')
        username = request.form.get('username')
        password = request.form.get('password')
        confirm_password = request.form.get('confirm_password')

        if not all([company_name, owner_name, email, phone, username, password, confirm_password]):
            return render_template('register_company.html', error='All fields are required')
        if password != confirm_password:
            return render_template('register_company.html', error='Passwords do not match')

        # Check if username already exists
        cred_ref = db.reference(f'credentials/users/{username}')
        if cred_ref.get():
            return render_template('register_company.html', error='Username already exists')

        # Validate phone number format
        clean_phone = re.sub(r'[^0-9]', '', phone)
        if len(clean_phone) < 10:
            return render_template('register_company.html', error='Please enter a valid phone number')

        # Store company data as "pending verification"
        temp_id = str(uuid.uuid4())
        hashed_pw = generate_password_hash(password)
        
        pending_data = {
            'temp_id': temp_id,
            'company_name': company_name,
            'owner_name': owner_name,
            'email': email,
            'phone': phone,
            'username': username,
            'password_hash': hashed_pw,
            'created_at': datetime.now().isoformat(),
            'phone_verified': False,
            'status': 'pending'  # FIX #1: Mark company as pending during registration
        }
        
        # Store in pending_companies
        db.reference(f'pending_companies/{temp_id}').set(pending_data)
        
        # Generate and send email OTP only
        email_otp = generate_otp()
        
        email_stored = store_email_otp_in_firebase(email, email_otp, username)
        
        if email_stored:
            # Send email OTP only
            send_email_otp(email, email_otp)
            
            # Redirect to email verification page only
            return render_template('verify_email.html', 
                                 temp_id=temp_id, 
                                 email=email,
                                 username=username,
                                 message='Please verify your email to complete registration. Check your inbox for verification code.')
        else:
            # Clean up pending data if OTP failed
            db.reference(f'pending_companies/{temp_id}').delete()
            return render_template('register_company.html', 
                                 error='Failed to send verification code. Please try again.')
    
    return render_template('register_company.html')

@app.route('/verify-registration', methods=['POST'])
def verify_registration():
    """Complete company registration after email verification only"""
    try:
        temp_id = request.form.get('temp_id')
        email_otp = request.form.get('email_otp')
        
        if not temp_id or not email_otp:
            return render_template('verify_email.html', 
                                 error='Email verification code is required')
        
        # Get pending company data
        pending_ref = db.reference(f'pending_companies/{temp_id}')
        pending_data = pending_ref.get()
        
        if not pending_data:
            return render_template('verify_email.html', 
                                 error='Registration session expired. Please register again.')
        
        email = pending_data.get('email')
        username = pending_data.get('username')
        
        # Verify email OTP only
        email_verification_success, email_message = verify_email_otp_from_firebase(email, email_otp)
        
        if email_verification_success:
            # Create actual user account
            cred_ref = db.reference(f'credentials/users/{username}')
            cred_ref.set({'password': pending_data['password_hash']})
            
            # Create user data with ACTIVE status
            db.reference(f'Data/{username}').set({
                'company_name': pending_data['company_name'],
                'owner_name': pending_data['owner_name'],
                'email': pending_data['email'],
                'phone': pending_data['phone'],
                'username': pending_data['username'],
                'password': pending_data['password_hash'],
                'phone_verified': False,  # Phone verification disabled
                'email_verified': True,
                'email_verified_at': datetime.now().isoformat(),
                'created_at': pending_data['created_at'],
                'status': 'active',  # FIX #2: Mark company as ACTIVE after email verification
                'activated_at': datetime.now().isoformat()
            })
            
            # Create storage folders for user
            db.reference(f'storage/{username}').set({
                'txt': {},
                'img': {},
                'audio': {},
                'files': {}
            })
            
            # Clean up pending data
            pending_ref.delete()
            
            return render_template('registration_success.html', 
                                 message='Registration completed successfully! Your email has been verified.')
        else:
            return render_template('verify_email.html', 
                                 temp_id=temp_id,
                                 email=email,
                                 username=username,
                                 error=email_message)
    
    except Exception as e:
        logger.error(f"Registration verification error: {str(e)}")
        return render_template('verify_email.html', 
                             error='An error occurred during verification. Please try again.')

@app.route('/resend-otp', methods=['POST'])
def resend_otp():
    """Resend OTP for registration verification - email only"""
    try:
        temp_id = request.form.get('temp_id', '').strip()
        
        # Validate temp_id
        if not temp_id:
            logger.warning("Resend OTP: temp_id is missing or empty")
            return jsonify({'success': False, 'error': 'Invalid request - session ID missing'})
        
        # Get pending company data from Firebase
        pending_ref = db.reference(f'pending_companies/{temp_id}')
        pending_data = pending_ref.get()
        
        if not pending_data:
            logger.warning(f"Resend OTP: No pending registration found for temp_id: {temp_id}")
            return jsonify({'success': False, 'error': 'Registration session expired. Please register again.'})
        
        email = pending_data.get('email')
        username = pending_data.get('username')
        
        if not email or not username:
            logger.error(f"Resend OTP: Missing email or username in pending data for temp_id: {temp_id}")
            return jsonify({'success': False, 'error': 'Invalid registration data'})
        
        # Generate new email OTP only
        email_otp = generate_otp()
        logger.info(f"Resend OTP: Generated new OTP for email: {email}")
        
        if store_email_otp_in_firebase(email, email_otp, username):
            send_email_otp(email, email_otp)
            logger.info(f"Resend OTP: Successfully sent OTP to {email}")
            
            # Check if in development mode to include dev_otp in response
            dev_otp = None
            try:
                otp_ref = db.reference(f'email_otps/{email}')
                otp_data = otp_ref.get()
                if otp_data and 'otp' in otp_data:
                    dev_otp = otp_data['otp']
            except:
                pass
            
            response = {'success': True, 'message': 'Email verification code sent successfully!'}
            if dev_otp:
                response['dev_email_otp'] = dev_otp
            
            return jsonify(response)
        else:
            logger.error(f"Resend OTP: Failed to store OTP in Firebase for email: {email}")
            return jsonify({'success': False, 'error': 'Failed to send email OTP'})
    
    except Exception as e:
        logger.error(f"Resend OTP error: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'error': 'Failed to resend OTP - server error'})

# Firebase Phone Auth Routes Removed - Email verification only

# Firebase SMS routes removed - Email verification only

@app.route('/login', methods=['GET', 'POST'])
def login():
    error = None
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        user_data = get_user_by_username(username)
        if not user_data:
            error = 'Invalid username'
        elif not check_password_hash(user_data['password'], password):
            error = 'Incorrect password'
        else:
            role = user_data.get('role', 'user')
            session['username'] = username
            session['role'] = role
            session['logged_in'] = True  # fix session persistence
            if role == 'admin':
                return redirect(url_for('admin_dashboard'))
            else:
                return redirect(url_for('index'))
    return render_template('login.html', error=error)

@app.route('/logout')
def logout():
    session.pop('logged_in', None)
    session.pop('username', None)
    session.pop('role', None)
    return redirect(url_for('login'))

# --- Admin Routes ---
@app.route('/admin/dashboard')
@admin_required
def admin_dashboard():
    return render_template('admin_dashboard.html', **get_user_context())

# --- Edit Company Endpoint ---
@app.route('/admin/edit-company', methods=['POST'])
@admin_required
def admin_edit_company():
    try:
        data = request.get_json()
        username = data.get('username')
        if not username:
            return jsonify({'success': False, 'error': 'Username required'})
        company_ref = db.reference(f'Data/{username}')
        company_data = company_ref.get()
        if not company_data:
            return jsonify({'success': False, 'error': 'Company not found'})
        # Update fields
        for field in ['company_name', 'owner_name', 'email', 'phone']:
            if field in data:
                company_data[field] = data[field]
        company_ref.set(company_data)
        return jsonify({'success': True})
    except Exception as e:
        logger.error(f"Edit company error: {str(e)}")
        return jsonify({'success': False, 'error': str(e)})

# --- Export Reports Endpoint ---
@app.route('/admin/export-reports', methods=['POST'])
@admin_required
def admin_export_reports():
    """Export admin reports as CSV, PDF, or JSON based on frontend request."""
    try:
        data = request.get_json()
        report_type = data.get('type', 'users')  # Add default value
        from_date = data.get('fromDate')
        to_date = data.get('toDate')
        export_format = data.get('format', 'csv')
        plan_filters = data.get('planFilters', [])

        # Normalize plan_filters to lowercase for comparison
        plan_filters_lower = [pf.lower() if isinstance(pf, str) else pf for pf in plan_filters]

        # Parse date range
        from_dt = datetime.strptime(from_date, '%Y-%m-%d') if from_date else None
        to_dt = datetime.strptime(to_date, '%Y-%m-%d') if to_date else None

        # Initialize result_data as empty list (critical: fresh for each request)
        result_data = []
        
        if report_type == 'users':
            users_ref = db.reference('Data')
            users_data = users_ref.get() or {}
            for username, user in users_data.items():
                if not isinstance(user, dict):
                    continue
                reg_date = user.get('registered_date', None)
                # Date filter (if available)
                if reg_date and from_dt and to_dt:
                    try:
                        reg_dt = datetime.strptime(reg_date, '%Y-%m-%d')
                        if not (from_dt <= reg_dt <= to_dt):
                            continue
                    except:
                        pass
                result_data.append({
                    'Username': username,
                    'Name': user.get('owner_name', ''),
                    'Email': user.get('email', ''),
                    'Phone': user.get('phone', ''),
                    'Company': user.get('company_name', ''),
                    'Plan': user.get('membership_status', ''),
                    'Registered': user.get('registered_date', '')
                })
        elif report_type == 'companies':
            companies_ref = db.reference('Data')
            companies_data = companies_ref.get() or {}
            for username, company in companies_data.items():
                if not isinstance(company, dict):
                    continue
                reg_date = company.get('registered_date', None)
                if reg_date and from_dt and to_dt:
                    try:
                        reg_dt = datetime.strptime(reg_date, '%Y-%m-%d')
                        if not (from_dt <= reg_dt <= to_dt):
                            continue
                    except:
                        pass
                result_data.append({
                    'Username': username,
                    'Company': company.get('company_name', ''),
                    'Owner': company.get('owner_name', ''),
                    'Email': company.get('email', ''),
                    'Phone': company.get('phone', ''),
                    'Plan': company.get('membership_status', ''),
                    'Registered': company.get('registered_date', '')
                })
        elif report_type == 'plans':
            users_ref = db.reference('Data')
            users_data = users_ref.get() or {}
            for username, user in users_data.items():
                if not isinstance(user, dict):
                    continue
                plan = user.get('membership_status', 'Free')
                # Case-insensitive plan filter comparison
                if plan_filters_lower and plan.lower() not in plan_filters_lower:
                    continue
                reg_date = user.get('registered_date', None)
                if reg_date and from_dt and to_dt:
                    try:
                        reg_dt = datetime.strptime(reg_date, '%Y-%m-%d')
                        if not (from_dt <= reg_dt <= to_dt):
                            continue
                    except:
                        pass
                result_data.append({
                    'Username': username,
                    'Plan': plan,
                    'Company': user.get('company_name', ''),
                    'Owner': user.get('owner_name', ''),
                    'Email': user.get('email', ''),
                    'Registered': user.get('registered_date', '')
                })
        elif report_type == 'usage':
            # Fetch usage analytics: tool usage and file processing stats from storage
            try:
                storage_ref = db.reference('storage')
                storage_data = storage_ref.get()
                
                # If storage doesn't exist, provide empty usage report
                if not storage_data or not isinstance(storage_data, dict):
                    # Fallback: Generate empty usage data for all registered users
                    users_ref = db.reference('Data')
                    users_data = users_ref.get() or {}
                    for username in users_data.keys():
                        result_data.append({
                            'Username': username,
                            'Total Files': 0,
                            'PDF Files': 0,
                            'Images': 0,
                            'Documents': 0,
                            'Usage Date': ''
                        })
                else:
                    # Process storage data by username
                    for username, user_storage in storage_data.items():
                        if not isinstance(user_storage, dict):
                            continue
                        
                        # Count files by category
                        total_files = 0
                        pdf_count = 0
                        image_count = 0
                        doc_count = 0
                        
                        for category, items in user_storage.items():
                            if isinstance(items, dict):
                                count = len(items)
                                total_files += count
                                if category.lower() == 'pdf':
                                    pdf_count = count
                                elif category.lower() in ['images', 'image']:
                                    image_count = count
                                elif category.lower() in ['documents', 'document', 'docs', 'doc']:
                                    doc_count = count
                        
                        result_data.append({
                            'Username': username,
                            'Total Files': total_files,
                            'PDF Files': pdf_count,
                            'Images': image_count,
                            'Documents': doc_count,
                            'Usage Date': datetime.now().strftime('%Y-%m-%d')
                        })
            except Exception as e:
                # If there's any error fetching storage, provide empty usage data
                logger.error(f"Error fetching usage data: {str(e)}")
                users_ref = db.reference('Data')
                users_data = users_ref.get() or {}
                for username in users_data.keys():
                    result_data.append({
                        'Username': username,
                        'Total Files': 0,
                        'PDF Files': 0,
                        'Images': 0,
                        'Documents': 0,
                        'Usage Date': ''
                    })
        elif report_type == 'financial':
            # Fetch financial data: users with their plan information
            users_ref = db.reference('Data')
            users_data = users_ref.get() or {}
            for username, user in users_data.items():
                if not isinstance(user, dict):
                    continue
                plan = user.get('membership_status', 'Free')
                if plan_filters_lower and plan.lower() not in plan_filters_lower:
                    continue
                reg_date = user.get('registered_date', None)
                if reg_date and from_dt and to_dt:
                    try:
                        reg_dt = datetime.strptime(reg_date, '%Y-%m-%d')
                        if not (from_dt <= reg_dt <= to_dt):
                            continue
                    except:
                        pass
                result_data.append({
                    'Username': username,
                    'Plan': plan,
                    'Company': user.get('company_name', ''),
                    'Owner': user.get('owner_name', ''),
                    'Email': user.get('email', ''),
                    'Registered': user.get('registered_date', ''),
                    'Status': 'Active' if user.get('membership_status') else 'Inactive'
                })
        elif report_type == 'comprehensive':
            # Combine all above
            # For brevity, just combine users and companies
            users_ref = db.reference('Data')
            users_data = users_ref.get() or {}
            for username, user in users_data.items():
                if not isinstance(user, dict):
                    continue
                result_data.append({
                    'Username': username,
                    'Name': user.get('owner_name', ''),
                    'Email': user.get('email', ''),
                    'Phone': user.get('phone', ''),
                    'Company': user.get('company_name', ''),
                    'Plan': user.get('membership_status', ''),
                    'Registered': user.get('registered_date', '')
                })

        # Output in requested format
        if export_format == 'json':
            return jsonify(result_data)
        elif export_format == 'csv':
            import csv
            output = io.StringIO()
            if result_data:
                writer = csv.DictWriter(output, fieldnames=result_data[0].keys())
                writer.writeheader()
                writer.writerows(result_data)
            else:
                output.write('No data found')
            output.seek(0)
            return send_file(
                io.BytesIO(output.getvalue().encode('utf-8')),
                mimetype='text/csv',
                as_attachment=True,
                download_name=f'docshift_{report_type}_report_{datetime.now().strftime("%Y%m%d")}.csv'
            )
        elif export_format == 'pdf':
            # Simple PDF export using FPDF
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font('Arial', 'B', 12)
            pdf.cell(0, 10, f'{report_type.title()} Report', ln=1, align='C')
            pdf.set_font('Arial', '', 10)
            if result_data:
                col_width = pdf.w / (len(result_data[0]) + 1)
                row_height = pdf.font_size * 1.5
                # Header
                for key in result_data[0].keys():
                    pdf.cell(col_width, row_height, str(key), border=1)
                pdf.ln(row_height)
                # Rows
                for row in result_data:
                    for val in row.values():
                        pdf.cell(col_width, row_height, str(val), border=1)
                    pdf.ln(row_height)
            else:
                pdf.cell(0, 10, 'No data found', ln=1)
            pdf_bytes = pdf.output(dest='S').encode('latin1')
            pdf_output = io.BytesIO(pdf_bytes)
            pdf_output.seek(0)
            return send_file(
                pdf_output,
                mimetype='application/pdf',
                as_attachment=True,
                download_name=f'docshift_{report_type}_report_{datetime.now().strftime("%Y%m%d")}.pdf'
            )
        else:
            return jsonify({'error': 'Invalid export format'}), 400
    except Exception as e:
        logger.error(f"Export reports error: {str(e)}")
        return jsonify({'error': f'Failed to export report: {str(e)}'}), 500

@app.route('/api/admin/dashboard-data')
@admin_required
def admin_dashboard_data():
    
    try:
        # ====== COMPANIES DATA FETCH ======
        companies_ref = db.reference('Data')
        companies_data = companies_ref.get() or {}
        
        # Process companies data
        companies_list = []
        total_companies = 0
        active_companies = 0
        
        for username, company_data in companies_data.items():
            if isinstance(company_data, dict) and 'company_name' in company_data:
                total_companies += 1
                
                # Get actual status from stored data, default to 'active' if missing
                company_status = company_data.get('status', 'active')
                
                if company_status == 'active':
                    active_companies += 1
                
                # Use actual status from database
                companies_list.append({
                    'company_name': company_data.get('company_name', ''),
                    'owner_name': company_data.get('owner_name', ''),
                    'email': company_data.get('email', ''),
                    'phone': company_data.get('phone', ''),
                    'username': username,
                    'status': company_status,
                    'registered_date': company_data.get('created_at', 'Recent')
                })
        
        # Sort by registered date (newest first)
        companies_list.sort(key=lambda x: x['registered_date'], reverse=True)
        
        # ====== USERS DATA FETCH ======
        users_list = []
        total_users = 0
        
        try:
            for username, user_data in companies_data.items():
                if isinstance(user_data, dict):
                    total_users += 1
                    users_list.append({
                        'id': username,
                        'username': username,
                        'email': user_data.get('email', ''),
                        'company': user_data.get('company_name', ''),
                        'plan': user_data.get('membership_status', 'Free')
                    })
            
            logger.info(f"Users fetched from Data reference: {total_users}")
            
        except Exception as e:
            logger.warning(f"Could not fetch users data: {str(e)}")
            users_list = []
            total_users = 0
        
        # Limit users to first 100
        users_list = users_list[:100]
        
        # ====== FILES/CONVERSIONS DATA FETCH ======
        files_list = []
        total_files = 0
        
        try:
            conversions_ref = db.reference('conversions')
            conversions_data = conversions_ref.get() or {}
            
            if isinstance(conversions_data, dict):
                for conversion_id, conversion_info in conversions_data.items():
                    if isinstance(conversion_info, dict):
                        files_list.append({
                            'id': conversion_id,
                            'conversion_type': conversion_info.get('type', 'unknown'),
                            'original_filename': conversion_info.get('original_name', ''),
                            'converted_filename': conversion_info.get('converted_name', ''),
                            'timestamp': conversion_info.get('timestamp', '')
                        })
                        total_files += 1
            
            logger.info(f"Files fetched from conversions: {len(files_list)}")
            
        except Exception as e:
            logger.warning(f"Could not fetch conversions data: {str(e)}")
            files_list = []
            total_files = 0
        
        # Fallback: If conversions is empty, try storage reference
        if total_files == 0:
            try:
                storage_ref = db.reference('storage')
                storage_data = storage_ref.get() or {}
                
                if isinstance(storage_data, dict):
                    for username_storage, user_storage in storage_data.items():
                        if isinstance(user_storage, dict):
                            for storage_type in ['txt', 'img', 'audio', 'files']:
                                type_data = user_storage.get(storage_type, {})
                                if isinstance(type_data, dict):
                                    for file_id, file_info in type_data.items():
                                        files_list.append({
                                            'id': file_id,
                                            'conversion_type': storage_type,
                                            'original_filename': file_info.get('name', file_id) if isinstance(file_info, dict) else file_id,
                                            'converted_filename': 'Stored',
                                            'timestamp': file_info.get('timestamp', '') if isinstance(file_info, dict) else ''
                                        })
                                        total_files += 1
                
                logger.info(f"Files built from storage: {len(files_list)}")
                
            except Exception as e:
                logger.warning(f"Could not fetch storage data: {str(e)}")
        
        # Limit files to last 100
        files_list = files_list[-100:] if len(files_list) > 100 else files_list
        
        # ====== BUILD RESPONSE - RETURN ALL COMPANIES ======
        dashboard_data = {
            'totalCompanies': total_companies,
            'activeCompanies': active_companies,
            'totalUsers': total_users,
            'totalFiles': total_files,
            'companies': companies_list,  # ✅ RETURN ALL COMPANIES (not just 10)
            'users': users_list,
            'files': files_list
        }
        
        return jsonify(dashboard_data)
        
    except Exception as e:
        logger.error(f"Error fetching admin dashboard data: {str(e)}")
        return jsonify({"error": "Failed to fetch data"}), 500

@app.route('/admin/companies')
@admin_required
def admin_companies():
    return render_template('admin_companies.html', **get_user_context())

@app.route('/api/admin/all-companies')
@admin_required
def admin_all_companies():
    try:
        # Get all company data from Firebase
        companies_ref = db.reference('Data')
        companies_data = companies_ref.get() or {}
        
        # Process companies data
        companies_list = []
        
        for username, company_data in companies_data.items():
            if isinstance(company_data, dict) and 'company_name' in company_data:
                companies_list.append({
                    'company_name': company_data.get('company_name', ''),
                    'owner_name': company_data.get('owner_name', ''),
                    'email': company_data.get('email', ''),
                    'phone': company_data.get('phone', ''),
                    'username': username,
                    'registered_date': 'Recent'  # You can add timestamp later
                })
        
        # Sort by company name
        companies_list.sort(key=lambda x: x['company_name'])
        
        return jsonify({'companies': companies_list})
        
    except Exception as e:
        logger.error(f"Error fetching all companies data: {str(e)}")
        return jsonify({"error": "Failed to fetch data"}), 500

@app.route('/admin/company-details')
@admin_required
def company_details():
    return render_template('company_details.html', **get_user_context())

@app.route('/api/admin/company-details/<username>')
@admin_required
def company_details_api(username):
    try:
        # Get company basic info
        company_ref = db.reference(f'Data/{username}')
        company_data = company_ref.get()
        
        if not company_data:
            return jsonify({"error": "Company not found"}), 404
        
        # Get storage data for usage statistics
        storage_ref = db.reference(f'storage/{username}')
        storage_data = storage_ref.get() or {}
        
        # Calculate statistics
        total_files = 0
        tool_usage = {
            'pdf': 0,
            'image': 0,
            'text': 0,
            'audio': 0,
            'compress': 0,
            'merge': 0,
            'split': 0,
            'convert': 0
        }
        
        # Count files by type
        for storage_type, files in storage_data.items():
            if isinstance(files, dict):
                file_count = len(files)
                total_files += file_count
                
                # Map storage types to tool usage
                if storage_type == 'txt':
                    tool_usage['text'] += file_count
                elif storage_type == 'img':
                    tool_usage['image'] += file_count
                elif storage_type == 'audio':
                    tool_usage['audio'] += file_count
                elif storage_type == 'files':
                    tool_usage['pdf'] += file_count
                    tool_usage['convert'] += file_count
        
        # Calculate approximate storage (assuming average file size)
        storage_used = total_files * 2.5  # Average 2.5 MB per file
        
        # Count unique tools used
        tools_used = sum(1 for count in tool_usage.values() if count > 0)
        
        # Generate sample recent activity
        recent_activity = []
        if total_files > 0:
            activities = [
                {"type": "pdf", "title": "PDF Converted", "description": "Word document converted to PDF", "time": "2 hours ago"},
                {"type": "image", "title": "Image Processed", "description": "Background removed from image", "time": "5 hours ago"},
                {"type": "text", "title": "Text Uploaded", "description": "Text content saved to storage", "time": "1 day ago"},
                {"type": "convert", "title": "File Conversion", "description": "Excel file converted to PDF", "time": "2 days ago"},
                {"type": "audio", "title": "Audio Generated", "description": "Text converted to speech", "time": "3 days ago"}
            ]
            recent_activity = activities[:min(5, total_files)]
        
        response_data = {
            'company': {
                'company_name': company_data.get('company_name', ''),
                'owner_name': company_data.get('owner_name', ''),
                'email': company_data.get('email', ''),
                'phone': company_data.get('phone', ''),
                'username': username
            },
            'stats': {
                'totalFiles': total_files,
                'storageUsed': round(storage_used, 1),
                'monthlyActivity': total_files,  # Simplified for now
                'toolsUsed': tools_used
            },
            'toolUsage': tool_usage,
            'recentActivity': recent_activity
        }
        
        return jsonify(response_data)
        
    except Exception as e:
        logger.error(f"Error fetching company details for {username}: {str(e)}")
        return jsonify({"error": "Failed to fetch company details"}), 500

@app.route('/admin/delete-company', methods=['POST'])
@admin_required
def delete_company():
    """Delete a company and all associated data"""
    try:
        # Get request data
        data = request.get_json()
        
        if not data:
            return jsonify({"error": "No data provided"}), 400
        
        username = data.get('username')
        confirm_delete = data.get('confirm_delete', False)
        
        if not username:
            return jsonify({"error": "Username is required"}), 400
        
        if not confirm_delete:
            return jsonify({"error": "Delete confirmation is required"}), 400
        
        # Check if company exists
        company_ref = db.reference(f'Data/{username}')
        company_data = company_ref.get()
        
        if not company_data:
            return jsonify({"error": "Company not found"}), 404
        
        # Get company name for logging
        company_name = company_data.get('company_name', username)
        
        # Log the deletion attempt
        logger.info(f"Admin attempting to delete company: {company_name} (username: {username}) by admin: {session.get('admin_id', 'unknown')}")
        
        # Start deletion process
        deletion_results = {
            'company_data': False,
            'storage_data': False,
            'credentials': False,
            'file_cleanup': False
        }
        
        try:
            # 1. Delete company basic data
            company_ref.delete()
            deletion_results['company_data'] = True
            logger.info(f"Deleted company data for: {username}")
            
        except Exception as e:
            logger.error(f"Error deleting company data for {username}: {str(e)}")
        
        try:
            # 2. Delete storage data (uploaded files metadata)
            storage_ref = db.reference(f'storage/{username}')
            storage_data = storage_ref.get()
            
            if storage_data:
                storage_ref.delete()
                deletion_results['storage_data'] = True
                logger.info(f"Deleted storage data for: {username}")
            else:
                deletion_results['storage_data'] = True  # No storage data to delete
                
        except Exception as e:
            logger.error(f"Error deleting storage data for {username}: {str(e)}")
        
        try:
            # 3. Delete user credentials if they exist
            credentials_ref = db.reference(f'credentials/{username}')
            credentials_data = credentials_ref.get()
            
            if credentials_data:
                credentials_ref.delete()
                deletion_results['credentials'] = True
                logger.info(f"Deleted credentials for: {username}")
            else:
                deletion_results['credentials'] = True  # No credentials to delete
                
        except Exception as e:
            logger.error(f"Error deleting credentials for {username}: {str(e)}")
        
        try:
            # 4. Clean up any physical files if they exist
            # This would depend on your file storage implementation
            # For now, we'll mark it as successful since files are typically auto-deleted
            deletion_results['file_cleanup'] = True
            
        except Exception as e:
            logger.error(f"Error during file cleanup for {username}: {str(e)}")
        
        # Check overall success
        total_operations = len(deletion_results)
        successful_operations = sum(deletion_results.values())
        
        if successful_operations == total_operations:
            # Complete success
            logger.info(f"Successfully deleted company: {company_name} (username: {username})")
            
            # Log to admin activity if you have such a system
            try:
                admin_activity_ref = db.reference('admin_activity')
                admin_activity_ref.push({
                    'action': 'delete_company',
                    'admin_id': session.get('admin_id', 'unknown'),
                    'target_company': company_name,
                    'target_username': username,
                    'timestamp': datetime.now().isoformat(),
                    'status': 'success'
                })
            except Exception as log_error:
                logger.error(f"Error logging admin activity: {str(log_error)}")
            
            return jsonify({
                "success": True,
                "message": f"Company '{company_name}' has been successfully deleted",
                "deletion_details": deletion_results
            })
            
        elif successful_operations > 0:
            # Partial success
            logger.warning(f"Partial deletion for company: {company_name} (username: {username}). Results: {deletion_results}")
            
            return jsonify({
                "success": True,
                "message": f"Company '{company_name}' has been partially deleted. Some data may remain.",
                "warning": "Partial deletion occurred",
                "deletion_details": deletion_results
            })
            
        else:
            # Complete failure
            logger.error(f"Failed to delete company: {company_name} (username: {username})")
            return jsonify({
                "error": "Failed to delete company. No data was removed.",
                "deletion_details": deletion_results
            }), 500
        
    except Exception as e:
        logger.error(f"Unexpected error during company deletion: {str(e)}")
        return jsonify({"error": f"Internal server error: {str(e)}"}), 500

@app.route('/admin/<path:filename>')
def admin_files(filename):
    """Serve admin static files"""
    return send_from_directory('admin', filename)

@app.route('/static/<path:filename>')
def static_files(filename):
    """Serve static files"""
    return send_from_directory('static', filename)

# Serve assests folder for FAQ images
@app.route('/assests/<path:filename>')
def assests_files(filename):
    return send_from_directory('assests', filename)

@app.route('/homepage')
def homepage():
    return render_template('homepage.html')

@app.route('/')
def index():
    if 'logged_in' not in session:
        return redirect(url_for('homepage'))
    return render_template('index.html', **get_user_context())

@app.route('/help')
@login_required
def help_support():
    """Help & Support page"""
    return render_template('help_support.html', **get_user_context())

@app.route('/privacy-policy')
def privacy_policy():
    """Privacy Policy page"""
    return render_template('privacy_policy.html', **get_user_context())

@app.route('/terms-conditions')
def terms_conditions():
    """Terms & Conditions page"""
    return render_template('terms_conditions.html', **get_user_context())

@app.route('/upgrade-plan')
@login_required
def upgrade_plan():
    """Upgrade Plan page"""
    return render_template('upgrade_plan.html', **get_user_context())

@app.route('/select_plan', methods=['POST'])
@login_required
def select_plan():
    """Handle plan selection and update user membership"""
    try:
        data = request.get_json()
        plan_name = data.get('plan', '')
        
        if not plan_name:
            return jsonify({'success': False, 'error': 'No plan selected'})
        
        # Validate plan name
        valid_plans = ['Free Trail', 'Standard', 'Premium']
        if plan_name not in valid_plans:
            return jsonify({'success': False, 'error': 'Invalid plan selected'})
        
        username = session.get('username')
        if not username:
            return jsonify({'success': False, 'error': 'User not authenticated'})
        
        # Update user's membership status in Firebase
        try:
            user_ref = db.reference(f'Data/{username}')
            user_data = user_ref.get() or {}
            user_data['membership_status'] = plan_name
            user_ref.set(user_data)
        except Exception as firebase_error:
            logger.warning(f"Firebase update failed: {str(firebase_error)}")
        
        # Update user's membership status in SQLite as backup
        try:
            conn = sqlite3.connect('file_conversion.db')
            cursor = conn.cursor()
            
            # Update the user's membership status
            cursor.execute("""
                UPDATE users 
                SET membership_status = ? 
                WHERE username = ?
            """, (plan_name, username))
            
            conn.commit()
            updated_rows = cursor.rowcount
            conn.close()
        except Exception as sqlite_error:
            logger.warning(f"SQLite update failed: {str(sqlite_error)}")
            updated_rows = 1  # Assume success if Firebase worked
        
        logger.info(f"User {username} upgraded to {plan_name} plan")
        return jsonify({
            'success': True, 
            'message': f'Successfully upgraded to {plan_name} plan!',
            'plan': plan_name
        })
            
    except Exception as e:
        logger.error(f"Plan selection error: {str(e)}")
        return jsonify({'success': False, 'error': 'An error occurred while selecting the plan'})

@app.route('/user_dashboard')
@login_required
def user_dashboard():
    return "User Dashboard (Under Construction)"

@app.route('/upload_txt', methods=['POST'])
@login_required
def upload_txt():
    content = request.form.get('text')
    if content and 'username' in session:
        user_path = get_user_storage_path()
        ref = db.reference(f'{user_path}/txt')
        ref.push({"content": content})
        return jsonify({"message": "Text uploaded successfully."})
    return jsonify({"error": "Invalid request"}), 400

# --- File Conversion and Upload/Download Routes (using per-user storage) ---

@app.route('/image-to-pdf')
@login_required
def image_to_pdf_page():
    return render_template('image_to_pdf.html', **get_user_context())

@app.route('/pdf-to-image')
@login_required
def pdf_to_image_page():
    return render_template('pdf_to_image.html', **get_user_context())

@app.route('/merge-pdfs')
@login_required
def merge_pdfs_page():
    return render_template('merge_pdfs.html', **get_user_context())

@app.route('/word-to-pdf')
@login_required
def word_to_pdf_page():
    return render_template('word_to_pdf.html', **get_user_context())

@app.route('/excel-to-pdf')
@login_required
def excel_to_pdf_page():
    return render_template('excel_to_pdf.html', **get_user_context())

@app.route('/pdf-to-ppt')
@login_required
def pdf_to_ppt_page():
    return render_template('pdf_to_ppt.html', **get_user_context())

@app.route('/bg-remover')
@login_required
def bg_remover_page():
    return render_template('background_remover.html', **get_user_context())

@app.route('/e-sign')
@login_required
def e_sign_page():
    return render_template('e_sign.html', **get_user_context())

@app.route('/admin-logs')
@login_required
def logs_page():
    return render_template('logs.html', **get_user_context())

@app.route('/history')
@login_required
def history_page():
    return render_template('history.html', **get_user_context())

@app.route('/compress-pdf')
@login_required
def compress_pdf_page():
    return render_template('compress_pdf.html', **get_user_context())

@app.route('/split-pdf')
@login_required
def split_pdf_page():
    return render_template('split_pdf.html', **get_user_context())

@app.route('/remove-pages-ui')
@login_required
def remove_pages_ui():
    return render_template('remove_page.html', **get_user_context())

@app.route('/document-screener')
@login_required
def document_screener_page():
    global current_document_text, conversation_history
    current_document_text = ''
    conversation_history.clear()
    return render_template('document_screener.html', **get_user_context())

@app.route('/plagiarism-scanner')
@login_required
def plagiarism_scanner_page():
    result = session.pop('plagiarism_result', None)
    input_text = session.pop('plagiarism_input_text', '')
    user_context = get_user_context()
    return render_template('plagiarism.html', result=result, input_text=input_text, **user_context)

@app.route('/text-to-speech')
@login_required
def text_to_speech_page():
    return render_template('text_to_speech.html', **get_user_context())

@app.route('/speech-to-text')
@login_required
def speech_to_text_page():
    return render_template('speech_to_text.html', **get_user_context())

@app.route('/ai-pdf-editor')
@login_required
def ai_pdf_editor_page():
    return render_template('ai_pdf_editor.html', **get_user_context())

@app.route('/text-summarizer')
@login_required
def text_summarizer_page():
    return render_template('text_summarizer.html', **get_user_context())

# --- Image to PDF ---
@app.route('/convert/image-to-pdf', methods=['POST'])
@login_required
def convert_image_to_pdf():
    if 'images' not in request.files:
        return jsonify({'error': 'No images provided'}), 400
    files = request.files.getlist('images')
    image_list = []
    for file in files:
        try:
            image = Image.open(file.stream)
            if image.mode != 'RGB':
                image = image.convert('RGB')
            image_list.append(image)
        except Exception as e:
            return jsonify({'error': f'Failed to read image: {str(e)}'}), 500
    if not image_list:
        return jsonify({'error': 'No valid images found'}), 400

    output_filename = "converted.pdf"
    image_id = uuid.uuid4().hex
    
    try:
        # Create temp directory for image conversions
        temp_dir = os.path.join(tempfile.gettempdir(), 'docshift_images')
        os.makedirs(temp_dir, exist_ok=True)
        temp_pdf_path = os.path.join(temp_dir, f'{image_id}.pdf')
        
        # Create PDF and save to temp file
        pdf_buffer = io.BytesIO()
        image_list[0].save(pdf_buffer, save_all=True, append_images=image_list[1:], format='PDF')
        pdf_buffer.seek(0)
        
        with open(temp_pdf_path, 'wb') as f:
            f.write(pdf_buffer.getvalue())
        
        # Store in session (metadata only, not binary)
        session['current_image_id'] = image_id
        session['current_image_path'] = temp_pdf_path
        session['current_image_filename'] = output_filename
        session.modified = True
        
        # Upload to cloudinary for permanent storage
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/files'
        
        pdf_buffer.seek(0)
        cloudinary_result = cloudinary.uploader.upload(pdf_buffer, 
                                                     folder=cloudinary_folder,
                                                     resource_type='raw',
                                                     public_id=output_filename)
        cloudinary_url = cloudinary_result['secure_url']
        store_url_in_firebase(cloudinary_url, 'files', output_filename)
        
        log_conversion('image-to-pdf', files[0].filename, output_filename, 'temp_file', cloudinary_url)
        
        # Return JSON with image_id (not file blob)
        file_size = os.path.getsize(temp_pdf_path)
        return jsonify({
            'status': 'success',
            'image_id': image_id,
            'filename': output_filename,
            'file_size': file_size
        }), 200
    except Exception as e:
        return jsonify({'error': f'PDF conversion failed: {str(e)}'}), 500

# --- Preview Image to PDF ---
@app.route('/preview-image-pdf')
@login_required
def preview_image_pdf():
    try:
        image_id = session.get('current_image_id')
        image_path = session.get('current_image_path')
        
        if not image_id or not image_path:
            return jsonify({'error': 'No conversion data in session'}), 400
        
        if not os.path.exists(image_path):
            return jsonify({'error': 'PDF file not found'}), 404
        
        # Display PDF in browser (as_attachment=False)
        return send_file(image_path, as_attachment=False, mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Preview failed: {str(e)}'}), 500

# --- Download Image to PDF ---
@app.route('/download-image-pdf')
@login_required
def download_image_pdf():
    try:
        image_id = session.get('current_image_id')
        image_path = session.get('current_image_path')
        filename = session.get('current_image_filename', 'converted.pdf')
        
        if not image_id or not image_path:
            return jsonify({'error': 'No conversion data in session'}), 400
        
        if not os.path.exists(image_path):
            return jsonify({'error': 'PDF file not found'}), 404
        
        # Download PDF (as_attachment=True)
        # Auto-cleanup temp file after download completes
        def cleanup_after_download():
            time.sleep(2)  # Wait 2 seconds to ensure download completes
            try:
                if os.path.exists(image_path):
                    os.remove(image_path)
                # Clear session data
                session.pop('current_image_id', None)
                session.pop('current_image_path', None)
                session.pop('current_image_filename', None)
            except Exception as e:
                print(f'Cleanup error: {str(e)}')
        
        # Start cleanup in background thread
        cleanup_thread = threading.Thread(target=cleanup_after_download, daemon=True)
        cleanup_thread.start()
        
        return send_file(image_path, as_attachment=True, download_name=filename, mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Download failed: {str(e)}'}), 500

# --- PDF to Image ---
@app.route('/convert/pdf-to-image', methods=['POST'])
@login_required
def convert_pdf_to_image():
    pdf_file = request.files.get('pdfFile')
    if not pdf_file:
        return jsonify({'error': 'No PDF uploaded'}), 400
    
    try:
        # Use PyMuPDF (fitz) to read PDF
        pdf_data = pdf_file.read()
        pdf_document = fitz.open(stream=pdf_data, filetype="pdf")
        
        if len(pdf_document) == 0:
            return jsonify({'error': 'PDF has no pages'}), 400
        
        # Create unique ID for this conversion
        pdf_id = uuid.uuid4().hex
        output_filename = "converted_images.zip"
        
        # Create temp directory for PDF image conversions
        temp_dir = os.path.join(tempfile.gettempdir(), 'docshift_pdf_images')
        os.makedirs(temp_dir, exist_ok=True)
        
        # Apply 300 DPI transformation for high-quality rendering (prevents pixelation at zoom)
        # Default PDF rendering is 72 DPI, so we scale by 300/72 ≈ 4.17x for crisp output
        dpi_scale = 300 / 72.0
        transform_matrix = fitz.Matrix(dpi_scale, dpi_scale)
        
        # Save first page as separate PNG for preview
        first_page = pdf_document.load_page(0)
        # Render at 300 DPI using transformation matrix for sharp preview
        first_pix = first_page.get_pixmap(matrix=transform_matrix)
        first_img_data = first_pix.tobytes("png")
        temp_preview_path = os.path.join(temp_dir, f'{pdf_id}_preview.png')
        with open(temp_preview_path, 'wb') as f:
            f.write(first_img_data)
        
        # Create ZIP with all pages
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                # Render each page at 300 DPI using transformation matrix for consistent high quality
                pix = page.get_pixmap(matrix=transform_matrix)
                img_data = pix.tobytes("png")
                zip_file.writestr(f'page_{page_num+1}.png', img_data)
        
        zip_buffer.seek(0)
        temp_zip_path = os.path.join(temp_dir, f'{pdf_id}.zip')
        
        # Save ZIP to temp file
        with open(temp_zip_path, 'wb') as f:
            f.write(zip_buffer.getvalue())
        
        # Store in session (metadata only, not binary)
        session['current_pdf_id'] = pdf_id
        session['current_pdf_path'] = temp_zip_path
        session['current_pdf_preview_path'] = temp_preview_path
        session['current_pdf_filename'] = output_filename
        session['pdf_page_count'] = len(pdf_document)
        session.modified = True
        
        # Upload to cloudinary for permanent storage
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/files'
        
        zip_buffer.seek(0)
        cloudinary_result = cloudinary.uploader.upload(zip_buffer, 
                                                     folder=cloudinary_folder,
                                                     resource_type='raw',
                                                     public_id=output_filename)
        cloudinary_url = cloudinary_result['secure_url']
        store_url_in_firebase(cloudinary_url, 'files', output_filename)
        
        log_conversion('pdf-to-image', pdf_file.filename, output_filename, 'temp_file', cloudinary_url)
        
        # Return JSON with pdf_id (not file blob)
        file_size = os.path.getsize(temp_zip_path)
        return jsonify({
            'status': 'success',
            'pdf_id': pdf_id,
            'filename': output_filename,
            'pages': len(pdf_document),
            'file_size': file_size
        }), 200
    except Exception as e:
        return jsonify({'error': f'Error during conversion: {str(e)}'}), 500

# --- Preview PDF to Image (First Page) ---
@app.route('/preview-pdf-image')
@login_required
def preview_pdf_image():
    try:
        pdf_id = session.get('current_pdf_id')
        preview_path = session.get('current_pdf_preview_path')
        
        if not pdf_id or not preview_path:
            return jsonify({'error': 'No conversion data in session'}), 400
        
        if not os.path.exists(preview_path):
            return jsonify({'error': 'Preview file not found'}), 404
        
        # Return PNG preview in browser
        return send_file(preview_path, as_attachment=False, mimetype='image/png')
    except Exception as e:
        return jsonify({'error': f'Preview failed: {str(e)}'}), 500

# --- Download PDF to Image (ZIP) ---
@app.route('/download-pdf-image')
@login_required
def download_pdf_image():
    try:
        pdf_id = session.get('current_pdf_id')
        pdf_path = session.get('current_pdf_path')
        preview_path = session.get('current_pdf_preview_path')
        filename = session.get('current_pdf_filename', 'converted_images.zip')
        
        if not pdf_id or not pdf_path:
            return jsonify({'error': 'No conversion data in session'}), 400
        
        if not os.path.exists(pdf_path):
            return jsonify({'error': 'ZIP file not found'}), 404
        
        # Download ZIP (as_attachment=True)
        # Auto-cleanup temp files after download completes
        def cleanup_after_download():
            time.sleep(2)  # Wait 2 seconds to ensure download completes
            try:
                if os.path.exists(pdf_path):
                    os.remove(pdf_path)
                if preview_path and os.path.exists(preview_path):
                    os.remove(preview_path)
                # Clear session data
                session.pop('current_pdf_id', None)
                session.pop('current_pdf_path', None)
                session.pop('current_pdf_preview_path', None)
                session.pop('current_pdf_filename', None)
                session.pop('pdf_page_count', None)
            except Exception as e:
                print(f'Cleanup error: {str(e)}')
        
        # Start cleanup in background thread
        cleanup_thread = threading.Thread(target=cleanup_after_download, daemon=True)
        cleanup_thread.start()
        
        return send_file(pdf_path, as_attachment=True, download_name=filename, mimetype='application/zip')
    except Exception as e:
        return jsonify({'error': f'Download failed: {str(e)}'}), 500


# --- Merge PDFs ---
@app.route('/merge/pdfs', methods=['POST'])
@login_required
def merge_pdfs():
    """
    Merge multiple PDFs and store temporarily for preview/download.
    Uses temp files instead of session to handle large binary data.
    """
    if 'pdfs' not in request.files:
        return jsonify({'error': 'No PDF files provided'}), 400
    
    files = request.files.getlist('pdfs')
    merger = PdfMerger()
    
    try:
        # Merge all PDFs
        for file in files:
            merger.append(file)
        
        output_filename = f"merged_{uuid.uuid4().hex}.pdf"
        
        # Create PDF in memory
        pdf_buffer = io.BytesIO()
        merger.write(pdf_buffer)
        merger.close()
        pdf_buffer.seek(0)
        
        # Get PDF data
        pdf_data = pdf_buffer.getvalue()
        
        # Create unique ID for this merged PDF
        merge_id = uuid.uuid4().hex
        
        # Save to temporary location on server
        temp_dir = os.path.join(tempfile.gettempdir(), 'docshift_merges')
        os.makedirs(temp_dir, exist_ok=True)
        temp_pdf_path = os.path.join(temp_dir, f'{merge_id}.pdf')
        
        with open(temp_pdf_path, 'wb') as f:
            f.write(pdf_data)
        
        # Store merge_id and filename in session (small data, no problem)
        session['current_merge_id'] = merge_id
        session['current_merge_filename'] = output_filename
        session['current_merge_path'] = temp_pdf_path
        session.modified = True
        
        # Count pages in merged PDF
        pdf_buffer.seek(0)
        pdf_reader = PdfReader(pdf_buffer)
        total_pages = len(pdf_reader.pages)
        
        # Upload to Cloudinary for permanent storage
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/files'
        
        pdf_buffer.seek(0)
        cloudinary_result = cloudinary.uploader.upload(pdf_buffer, 
                                                     folder=cloudinary_folder,
                                                     resource_type='raw',
                                                     public_id=output_filename)
        cloudinary_url = cloudinary_result['secure_url']
        store_url_in_firebase(cloudinary_url, 'files', output_filename)
        
        log_conversion('merge-pdfs', files[0].filename, output_filename, 'memory', cloudinary_url)
        
        # Return JSON response with merge ID (NOT large binary data)
        return jsonify({
            'status': 'success',
            'message': 'PDFs merged successfully!',
            'merge_id': merge_id,
            'filename': 'merged.pdf',
            'pages': total_pages,
            'file_size': len(pdf_data)
        }), 200
        
    except Exception as e:
        logger.error(f"Merge error: {str(e)}")
        return jsonify({'error': f'Merge failed: {str(e)}'}), 500
    except Exception as e:
        return jsonify({'error': f'Merge failed: {str(e)}'}), 500

# --- Download Merged PDF ---
@app.route('/download-merged', methods=['GET'])
@login_required
def download_merged():
    """
    Download the previously merged PDF.
    Called after user previews and confirms download.
    """
    merge_id = session.get('current_merge_id')
    merge_path = session.get('current_merge_path')
    filename = session.get('current_merge_filename', 'merged.pdf')
    
    if not merge_id or not merge_path:
        return jsonify({'error': 'No merged PDF in session. Please merge PDFs first.'}), 400
    
    try:
        # Check if file exists
        if not os.path.exists(merge_path):
            return jsonify({'error': 'Merged PDF file not found. Please merge PDFs again.'}), 404
        
        # Return file for download
        response = send_file(
            merge_path,
            as_attachment=True,
            download_name=filename,
            mimetype='application/pdf'
        )
        
        # Clean up after successful download (in background, don't wait)
        def cleanup():
            import time
            time.sleep(2)  # Wait 2 seconds to ensure download completes
            try:
                if os.path.exists(merge_path):
                    os.remove(merge_path)
                # Clear session
                session.pop('current_merge_id', None)
                session.pop('current_merge_path', None)
                session.pop('current_merge_filename', None)
            except:
                pass
        
        thread = threading.Thread(target=cleanup, daemon=True)
        thread.start()
        
        return response
        
    except Exception as e:
        logger.error(f"Download error: {str(e)}")
        return jsonify({'error': f'Download failed: {str(e)}'}), 500

# --- Preview Merged PDF ---
@app.route('/preview-merged', methods=['GET'])
@login_required
def preview_merged():
    """
    Preview the merged PDF in browser (new tab).
    Retrieves from temporary file storage.
    """
    merge_id = session.get('current_merge_id')
    merge_path = session.get('current_merge_path')
    
    if not merge_id or not merge_path:
        return jsonify({'error': 'No merged PDF in session. Please merge PDFs first.'}), 400
    
    try:
        # Check if file exists
        if not os.path.exists(merge_path):
            return jsonify({'error': 'Merged PDF file not found. Please merge PDFs again.'}), 404
        
        # Return file WITHOUT as_attachment=True (displays in browser instead of downloading)
        return send_file(
            merge_path,
            as_attachment=False,
            mimetype='application/pdf'
        )
    except Exception as e:
        logger.error(f"Preview error: {str(e)}")
        return jsonify({'error': f'Preview failed: {str(e)}'}), 500

# --- Download route (per-user) ---
@app.route('/download')
@login_required
def download():
    file_path = request.args.get('file_path')
    file_name = request.args.get('file_name')
    mime_type = request.args.get('mime_type', 'application/octet-stream')
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found on server'}), 404

    ext = os.path.splitext(file_name)[1].lower()
    if ext in ['.pdf', '.docx', '.xlsx', '.pptx', '.txt', '.csv']:
        category = 'files'
    elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp']:
        category = 'img'
    elif ext in ['.mp3', '.wav', '.aac', '.ogg', '.flac']:
        category = 'audio'
    else:
        category = 'files'

    username = session.get('username')
    folder_path = f'storage/{username}/{category}'
    
    # Read file into memory for upload and sending
    with open(file_path, 'rb') as f:
        file_data = f.read()
    file_buffer = io.BytesIO(file_data)
    
    # Upload from memory
    file_buffer.seek(0)
    cloudinary_result = cloudinary.uploader.upload(file_buffer, 
                                                 folder=folder_path,
                                                 resource_type='raw',
                                                 public_id=file_name)
    cloudinary_url = cloudinary_result['secure_url']
    store_url_in_firebase(cloudinary_url, category, file_name)

    # Send from memory
    file_buffer.seek(0)
    return send_file(file_buffer, as_attachment=True, download_name=file_name, mimetype=mime_type)

# --- Download from Cloudinary/Firebase ---
@app.route('/download-file/<int:file_id>')
@login_required
def download_file_from_cloud(file_id):
    """Download file from Cloudinary using file ID from database."""
    try:
        # Get file info from database
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        c.execute('''
            SELECT cloudinary_url, converted_filename, original_filename 
            FROM conversions 
            WHERE id = ? AND username = ?
        ''', (file_id, session.get('username')))
        result = c.fetchone()
        conn.close()
        
        if not result:
            return jsonify({'error': 'File not found or access denied'}), 404
            
        cloudinary_url, converted_filename, original_filename = result
        
        if not cloudinary_url:
            return jsonify({'error': 'File URL not available'}), 404
        
        # Fetch file from Cloudinary
        response = requests.get(cloudinary_url)
        if response.status_code != 200:
            return jsonify({'error': 'File not accessible from cloud storage'}), 404
        
        # Create a temporary file to serve
        temp_file = tempfile.NamedTemporaryFile(delete=False)
        temp_file.write(response.content)
        temp_file.close()
        
        # Determine MIME type
        mime_type, _ = mimetypes.guess_type(converted_filename)
        if not mime_type:
            mime_type = 'application/octet-stream'
        
        def cleanup_temp_file():
            try:
                os.unlink(temp_file.name)
            except:
                pass
        
        # Schedule cleanup after sending file
        from threading import Timer
        Timer(5.0, cleanup_temp_file).start()
        
        return send_file(
            temp_file.name, 
            as_attachment=True, 
            download_name=converted_filename, 
            mimetype=mime_type
        )
        
    except Exception as e:
        return jsonify({'error': f'Download failed: {str(e)}'}), 500

# --- Word to PDF ---
@app.route('/word_to_pdf', methods=['POST'])
@login_required
def convert_word_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400
    file = request.files['file']

    if not file or not file.filename:
        return jsonify({'error': 'No file selected!'}), 400

    filename = file.filename.lower()
    original_filename = file.filename

    allowed_extensions = ['.doc', '.docx']
    file_extension = '.' + filename.split('.')[-1] if '.' in filename else ''

    if file_extension not in allowed_extensions:
        return jsonify({'error': 'Only .doc and .docx files are supported'}), 400

    # Save uploaded file to a temporary file and close immediately
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_in:
        file.save(temp_in)
        input_path = temp_in.name

    # Use temporary file for output in docshift_word_pdf directory
    word_id = uuid.uuid4().hex
    temp_dir = os.path.join(tempfile.gettempdir(), 'docshift_word_pdf')
    os.makedirs(temp_dir, exist_ok=True)
    output_path = os.path.join(temp_dir, f'{word_id}.pdf')

    try:
        if file_extension == '.docx':
            try:
                import pythoncom
                pythoncom.CoInitialize()
                from docx2pdf import convert as docx2pdf_convert
                docx2pdf_convert(input_path, output_path)
                pythoncom.CoUninitialize()
                if not os.path.exists(output_path):
                    raise Exception("docx2pdf did not create output file.")
            except Exception as e:
                logger.error(f'docx2pdf failed: {e}, falling back to Platypus PDF conversion.')
                # --- Platypus fallback: full content, proper wrapping, tables ---
                from reportlab.lib.units import cm as _cm

                def _safe(txt):
                    return txt.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').translate(
                        {i: '?' for i in range(65536, 1114112)})

                _styles = getSampleStyleSheet()
                _normal = ParagraphStyle('_N', parent=_styles['Normal'], fontSize=11, leading=15, spaceAfter=5)
                _h1 = ParagraphStyle('_H1', parent=_styles['Normal'], fontSize=14, leading=18, spaceBefore=10, spaceAfter=6, fontName='Helvetica-Bold')
                _h2 = ParagraphStyle('_H2', parent=_styles['Normal'], fontSize=12, leading=16, spaceBefore=8, spaceAfter=4, fontName='Helvetica-Bold')

                doc_in = Document(input_path)
                story = []

                for para in doc_in.paragraphs:
                    raw = para.text
                    if not raw.strip():
                        story.append(Spacer(1, 4))
                        continue
                    sname = para.style.name if para.style else ''
                    if 'Heading 1' in sname or 'Title' in sname:
                        story.append(Paragraph(_safe(raw), _h1))
                    elif 'Heading' in sname:
                        story.append(Paragraph(_safe(raw), _h2))
                    else:
                        story.append(Paragraph(_safe(raw), _normal))

                for tbl in doc_in.tables:
                    tdata = []
                    for row in tbl.rows:
                        tdata.append([Paragraph(_safe(c.text.strip()), _normal) for c in row.cells])
                    if tdata:
                        ncols = max(len(r) for r in tdata)
                        cw = (A4[0] - 4 * _cm) / max(ncols, 1)
                        t = Table(tdata, colWidths=[cw] * ncols)
                        t.setStyle(TableStyle([
                            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                            ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                            ('FONTSIZE', (0, 0), (-1, -1), 10),
                            ('LEFTPADDING', (0, 0), (-1, -1), 4),
                            ('RIGHTPADDING', (0, 0), (-1, -1), 4),
                        ]))
                        story.append(t)
                        story.append(Spacer(1, 8))

                if not story:
                    raise Exception("No content found in docx for fallback conversion.")

                pdf_io = io.BytesIO()
                SimpleDocTemplate(pdf_io, pagesize=A4,
                                  leftMargin=2*_cm, rightMargin=2*_cm,
                                  topMargin=2*_cm, bottomMargin=2*_cm).build(story)
                pdf_io.seek(0)
                with open(output_path, 'wb') as f:
                    f.write(pdf_io.getvalue())

        elif file_extension == '.doc':
            # Handle legacy .doc files
            try:
                # Attempt to open .doc as .docx (may fail)
                doc = Document(input_path)
                text = [para.text for para in doc.paragraphs if para.text.strip()]

                if not text:
                    raise Exception("No text extracted from .doc file using python-docx")

            except Exception as e:
                logger.warning(f"python-docx failed for .doc file: {e}, trying alternative method")
                try:
                    with open(input_path, 'rb') as f:
                        content = f.read()
                    text_content = content.decode('utf-8', errors='ignore')
                    lines = text_content.split('\n')
                    text = [line.strip() for line in lines if line.strip() and len(line.strip()) > 2]

                    clean_text = []
                    for line in text[:100]:
                        if any(c.isalpha() for c in line) and len(line) < 200:
                            clean_text.append(line)
                    if not clean_text:
                        raise Exception("No readable text found in .doc file")

                    text = clean_text

                except Exception as e2:
                    logger.error(f"Alternative .doc reading method failed: {e2}")
                    return jsonify({'error': 'Failed to read .doc file. Please try converting to .docx format first.'}), 500

            pdf_io = io.BytesIO()
            from reportlab.lib.units import cm as _cm2

            def _safe2(txt):
                return txt.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').translate(
                    {i: '?' for i in range(65536, 1114112)})

            _st = getSampleStyleSheet()
            _np = ParagraphStyle('_NP', parent=_st['Normal'], fontSize=11, leading=15, spaceAfter=5)
            story2 = [Paragraph(_safe2(ln.strip()), _np) for ln in text if ln.strip()]

            if not story2:
                raise Exception("No content found in .doc file for PDF conversion.")

            SimpleDocTemplate(pdf_io, pagesize=A4,
                              leftMargin=2*_cm2, rightMargin=2*_cm2,
                              topMargin=2*_cm2, bottomMargin=2*_cm2).build(story2)
            pdf_io.seek(0)
            with open(output_path, 'wb') as f:
                f.write(pdf_io.getvalue())

        # Verify PDF creation
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            raise Exception("PDF conversion failed - output file is empty or missing")

        # Store in session for preview/download endpoints
        session['word_result_id'] = word_id
        session['word_result_path'] = output_path
        session['word_result_filename'] = original_filename
        session.modified = True
        
        # Return JSON response (no auto-download - user chooses preview or download)
        file_size = os.path.getsize(output_path)
        return jsonify({
            'status': 'success',
            'word_id': word_id,
            'filename': original_filename,
            'file_size': file_size
        }), 200

    except Exception as e:
        logger.error(f'Word to PDF conversion failed: {str(e)}', exc_info=True)
        return jsonify({'error': f'Word to PDF conversion failed: {str(e)}'}), 500

    finally:
        # Clean up temp input file
        try:
            if os.path.exists(input_path):
                os.remove(input_path)
        except Exception as cleanup_error:
            logger.warning(f"Error deleting temp input file {input_path}: {str(cleanup_error)}")

@app.route('/preview-word-pdf', methods=['GET'])
@login_required
def preview_word_pdf():
    result_path = session.get('word_result_path')
    if not result_path or not os.path.exists(result_path):
        return "File not found", 404
    
    return send_file(result_path, mimetype='application/pdf', as_attachment=False)

@app.route('/download-word-pdf', methods=['GET'])
@login_required
def download_word_pdf():
    result_path = session.get('word_result_path')
    result_filename = session.get('word_result_filename', 'converted.pdf')
    
    if not result_path or not os.path.exists(result_path):
        return "File not found", 404
    
    # Log conversion before download
    log_conversion('word-to-pdf', result_filename, 'converted.pdf', 'file', result_path, 'success')
    
    # Upload to Cloudinary
    username = session.get('username')
    cloudinary_folder = f'storage/{username}/files'
    
    try:
        cloudinary_result = cloudinary.uploader.upload(result_path, 
                                                      folder=cloudinary_folder,
                                                      resource_type='raw',
                                                      public_id=f'word_{uuid.uuid4().hex}.pdf')
        cloudinary_url = cloudinary_result['secure_url']
        store_url_in_firebase(cloudinary_url, 'files', 'converted.pdf')
    except Exception as e:
        print(f"Cloudinary upload error: {e}")
    
    # Background cleanup
    def cleanup():
        import time
        time.sleep(2)
        if os.path.exists(result_path):
            os.remove(result_path)
        session.pop('word_result_id', None)
        session.pop('word_result_path', None)
        session.pop('word_result_filename', None)
    
    cleanup_thread = threading.Thread(target=cleanup, daemon=True)
    cleanup_thread.start()
    
    base_name = result_filename.rsplit('.', 1)[0] if '.' in result_filename else result_filename
    download_name = f"{base_name}.pdf"
    
    return send_file(result_path, as_attachment=True, download_name=download_name, mimetype='application/pdf')

# --- Excel to PDF ---
@app.route('/convert_excel_to_pdf', methods=['POST'])
@login_required
def convert_excel_to_pdf():
    # 1. Validate file upload
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400

    file = request.files['file']

    if not (file.filename.endswith('.xlsx') or file.filename.endswith('.xls') or file.filename.endswith('.csv')):
        return jsonify({'error': 'Only .xlsx, .xls, and .csv files are supported'}), 400

    # 2. Save Excel file temporarily
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_input:
        file.save(temp_input)
        input_path = temp_input.name

    output_filename = "converted.pdf"
    excel_id = uuid.uuid4().hex

    # 3. Create temporary output PDF file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_output:
        output_path = temp_output.name

    try:
        # 4. Read Excel/CSV data and convert to structured table format
        try:
            all_tables = []  # Store all table data for each sheet/CSV
            
            if file.filename.endswith('.csv'):
                # Handle CSV files - convert to table structure
                table_data = []
                with open(input_path, 'r', encoding='utf-8') as csvfile:
                    csv_reader = csv.reader(csvfile)
                    for row in csv_reader:
                        # Convert all cells to strings for table rendering
                        row_data = [str(cell) if cell else '' for cell in row]
                        if any(cell.strip() for cell in row_data):  # Only add non-empty rows
                            table_data.append(row_data)
                
                if table_data:
                    all_tables.append(("CSV Data", table_data))
                else:
                    all_tables.append(("CSV Data", [["(No data in CSV file)"]]))
            else:
                # Handle Excel files - convert each sheet to table structure
                wb = load_workbook(input_path)
                if not wb.sheetnames:
                    return jsonify({'error': 'Excel file contains no worksheets'}), 400

                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    table_data = []
                    
                    # Extract all rows from the sheet
                    for row in sheet.rows:
                        row_data = [
                            str(cell.value) if cell.value is not None else ''
                            for cell in row
                        ]
                        if any(cell.strip() for cell in row_data if cell):  # Only add non-empty rows
                            table_data.append(row_data)
                    
                    if not table_data:
                        table_data = [["(No data in this sheet)"]]
                    
                    all_tables.append((sheet_name, table_data))
        except Exception as file_error:
            return jsonify({'error': f'Invalid file: {str(file_error)}'}), 400

        # 5. Generate PDF with Platypus Table rendering (structured columns instead of flat text)
        # Using Platypus Table ensures proper tabular formatting, preserves column alignment,
        # and handles multi-page tables correctly without breaking data structure
        pdf_io = io.BytesIO()
        doc = SimpleDocTemplate(pdf_io, pagesize=letter, 
                              topMargin=0.5*inch, bottomMargin=0.5*inch,
                              leftMargin=0.5*inch, rightMargin=0.5*inch)
        
        story = []  # Content to render
        styles = getSampleStyleSheet()
        
        # Define table styling for professional appearance
        table_style = TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),  # Header row background
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),  # Header text color
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),  # Left align all cells
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),  # Bold header
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('FONTSIZE', (0, 1), (-1, -1), 9),  # Normal data font size
            ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),  # Grid lines for clarity
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.beige, colors.white]),  # Alternating row colors
        ])
        
        # Add each sheet/CSV as a separate table
        for sheet_name, table_data in all_tables:
            # Add sheet name as title
            title = Paragraph(f"<b>{sheet_name}</b>", styles['Heading3'])
            story.append(title)
            story.append(Spacer(1, 0.2*inch))
            
            # Create table from structured data (rows and columns preserved)
            table = Table(table_data, repeatRows=1)  # Repeat header row on page breaks
            table.setStyle(table_style)
            story.append(table)
            story.append(Spacer(1, 0.3*inch))  # Space between sheets
        
        # Build PDF document with all content
        doc.build(story)
        pdf_io.seek(0)

        # 6. Write PDF to temp output file
        with open(output_path, 'wb') as f:
            f.write(pdf_io.getvalue())

        # 7. Create temp directory for Excel conversions
        temp_dir = os.path.join(tempfile.gettempdir(), 'docshift_excel')
        os.makedirs(temp_dir, exist_ok=True)
        temp_pdf_path = os.path.join(temp_dir, f'{excel_id}.pdf')
        
        # Copy to temp dir
        with open(output_path, 'rb') as src:
            with open(temp_pdf_path, 'wb') as dst:
                dst.write(src.read())
        
        # 8. Store in session (metadata only, not binary)
        session['current_excel_id'] = excel_id
        session['current_excel_path'] = temp_pdf_path
        session['current_excel_filename'] = output_filename
        session.modified = True
        
        # 9. Upload to cloudinary for permanent storage
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/files'
        
        pdf_io.seek(0)
        cloudinary_result = cloudinary.uploader.upload(pdf_io, 
                                                     folder=cloudinary_folder,
                                                     resource_type='raw',
                                                     public_id=output_filename)
        cloudinary_url = cloudinary_result['secure_url']
        store_url_in_firebase(cloudinary_url, 'files', output_filename)
        
        log_conversion('excel-to-pdf', file.filename, output_filename, 'temp_file', cloudinary_url)
        
        # 10. Return JSON with excel_id (not file blob)
        file_size = os.path.getsize(temp_pdf_path)
        return jsonify({
            'status': 'success',
            'excel_id': excel_id,
            'filename': output_filename,
            'file_size': file_size
        }), 200

    except Exception as e:
        return jsonify({
            'error': f'Excel to PDF conversion failed: {str(e)}'
        }), 500

    finally:
        # Clean up temporary files
        if os.path.exists(input_path):
            os.remove(input_path)
        if 'output_path' in locals() and os.path.exists(output_path):
            os.remove(output_path)

# --- Preview Excel to PDF ---
@app.route('/preview-excel-pdf')
@login_required
def preview_excel_pdf():
    try:
        excel_id = session.get('current_excel_id')
        excel_path = session.get('current_excel_path')
        
        if not excel_id or not excel_path:
            return jsonify({'error': 'No conversion data in session'}), 400
        
        if not os.path.exists(excel_path):
            return jsonify({'error': 'PDF file not found'}), 404
        
        # Display PDF in browser (as_attachment=False)
        return send_file(excel_path, as_attachment=False, mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Preview failed: {str(e)}'}), 500

# --- Download Excel to PDF ---
@app.route('/download-excel-pdf')
@login_required
def download_excel_pdf():
    try:
        excel_id = session.get('current_excel_id')
        excel_path = session.get('current_excel_path')
        filename = session.get('current_excel_filename', 'converted.pdf')
        
        if not excel_id or not excel_path:
            return jsonify({'error': 'No conversion data in session'}), 400
        
        if not os.path.exists(excel_path):
            return jsonify({'error': 'PDF file not found'}), 404
        
        # Download PDF (as_attachment=True)
        # Auto-cleanup temp file after download completes
        def cleanup_after_download():
            time.sleep(2)  # Wait 2 seconds to ensure download completes
            try:
                if os.path.exists(excel_path):
                    os.remove(excel_path)
                # Clear session data
                session.pop('current_excel_id', None)
                session.pop('current_excel_path', None)
                session.pop('current_excel_filename', None)
            except Exception as e:
                print(f'Cleanup error: {str(e)}')
        
        # Start cleanup in background thread
        cleanup_thread = threading.Thread(target=cleanup_after_download, daemon=True)
        cleanup_thread.start()
        
        return send_file(excel_path, as_attachment=True, download_name=filename, mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Download failed: {str(e)}'}), 500



# --- PDF to PPT ---
@app.route('/pdf_to_ppt', methods=['POST'])
@login_required
def convert_pdf_to_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400
    file = request.files['file']
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'Only .pdf files are supported'}), 400

    original_filename = file.filename
    ppt_id = uuid.uuid4().hex
    temp_dir = os.path.join(tempfile.gettempdir(), 'docshift_ppt')
    os.makedirs(temp_dir, exist_ok=True)
    output_path = os.path.join(temp_dir, f'{ppt_id}.pptx')

    try:
        from pptx.dml.color import RGBColor
        from pptx.util import Emu

        file.stream.seek(0)
        pdf_bytes = file.stream.read()
        pdf_doc = fitz.open(stream=pdf_bytes, filetype='pdf')

        prs = Presentation()
        # Match slide size to PDF page dimensions of first page (or default 16:9)
        first_page = pdf_doc[0]
        pdf_w = first_page.rect.width   # points
        pdf_h = first_page.rect.height  # points
        # Convert PDF points → EMU (1 pt = 12700 EMU)
        prs.slide_width  = Emu(int(pdf_w * 12700))
        prs.slide_height = Emu(int(pdf_h * 12700))

        blank_layout = prs.slide_layouts[6]  # fully blank — no placeholders

        # Render scale: 2x for crisp quality without being too heavy
        RENDER_SCALE = 2.0
        mat = fitz.Matrix(RENDER_SCALE, RENDER_SCALE)

        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            page_w = page.rect.width
            page_h = page.rect.height

            # ── 1. Render full page as PNG → slide background ────────────
            # This preserves ALL visuals: background color, gradients,
            # decorative shapes, images, rendered fonts — pixel-perfect.
            pixmap = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pixmap.tobytes("png")
            bg_buf = io.BytesIO(img_bytes)
            bg_buf.seek(0)

            slide = prs.slides.add_slide(blank_layout)
            # Place background image covering the entire slide
            slide.shapes.add_picture(
                bg_buf,
                left=0, top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )

            # ── 2. Overlay transparent text boxes with positioned text ───
            # Extract rich text dict: preserves per-span font, size, color,
            # bold, italic, and exact bounding box on the page.
            text_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)

            for block in text_dict.get("blocks", []):
                if block.get("type") != 0:   # 0 = text block; skip image blocks
                    continue

                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        span_text = span.get("text", "").strip()
                        if not span_text:
                            continue

                        # Bounding box of this span in PDF points
                        x0, y0, x1, y1 = span["bbox"]
                        span_w = max(x1 - x0, 1)
                        span_h = max(y1 - y0, 1)

                        # Scale from PDF points → EMU
                        left   = Emu(int(x0 * 12700))
                        top    = Emu(int(y0 * 12700))
                        width  = Emu(int(span_w * 12700))
                        height = Emu(int(span_h * 12700 * 1.4))  # slight pad so text fits

                        txb = slide.shapes.add_textbox(left, top, width, height)
                        tf  = txb.text_frame
                        tf.word_wrap = False

                        p   = tf.paragraphs[0]
                        run = p.add_run()
                        run.text = span_text

                        # Font size from PDF (in points)
                        font_size = span.get("size", 12)
                        run.font.size = Pt(max(font_size, 6))

                        # Bold / italic from PDF flags
                        flags = span.get("flags", 0)
                        run.font.bold   = bool(flags & 2 ** 4)   # bit 4 = bold
                        run.font.italic = bool(flags & 2 ** 1)   # bit 1 = italic

                        # Font color from PDF (stored as 0xRRGGBB int)
                        raw_color = span.get("color", 0)
                        if raw_color is not None:
                            try:
                                r = (raw_color >> 16) & 0xFF
                                g = (raw_color >> 8)  & 0xFF
                                b =  raw_color        & 0xFF
                                run.font.color.rgb = RGBColor(r, g, b)
                            except Exception:
                                pass  # leave color default if anything goes wrong

                        # Font name from PDF when available
                        font_name = span.get("font", "")
                        if font_name:
                            # Clean up PDF font name suffixes like "ABCDEF+Arial-Bold"
                            clean_name = font_name.split("+")[-1].split("-")[0].strip()
                            if clean_name:
                                try:
                                    run.font.name = clean_name
                                except Exception:
                                    pass

                        # Make textbox background transparent (no fill)
                        from pptx.oxml.ns import qn
                        from lxml import etree
                        spPr = txb._element.find(qn('p:sp') + '/' + qn('p:spPr'), txb._element.nsmap)
                        if spPr is None:
                            spPr = txb._element.spPr if hasattr(txb._element, 'spPr') else None
                        if spPr is not None:
                            # Remove any existing fill and set noFill
                            for existing in spPr.findall(qn('a:solidFill')):
                                spPr.remove(existing)
                            no_fill = etree.SubElement(spPr, qn('a:noFill'))  # noqa

        pdf_doc.close()
        prs.save(output_path)

        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            raise Exception("PPT conversion failed - output file is empty or missing")

        session['ppt_result_id']       = ppt_id
        session['ppt_result_path']     = output_path
        session['ppt_result_filename'] = original_filename
        session.modified = True

        return jsonify({
            'status': 'success',
            'ppt_id': ppt_id,
            'filename': original_filename,
            'file_size': os.path.getsize(output_path)
        }), 200

    except Exception as e:
        logger.error(f'PDF to PPT conversion failed: {str(e)}', exc_info=True)
        return jsonify({'error': f'PDF to PPT conversion failed: {str(e)}'}), 500

@app.route('/preview-ppt', methods=['GET'])
@login_required
def preview_ppt():
    result_path = session.get('ppt_result_path')
    if not result_path or not os.path.exists(result_path):
        return "File not found", 404
    
    base_name = session.get('ppt_result_filename', 'presentation').rsplit('.', 1)[0]
    download_name = f"{base_name}.pptx"
    
    # Open in new tab without forcing download
    return send_file(result_path, as_attachment=False, download_name=download_name, 
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

@app.route('/download-ppt', methods=['GET'])
@login_required
def download_ppt():
    result_path = session.get('ppt_result_path')
    result_filename = session.get('ppt_result_filename', 'presentation.pdf')
    
    if not result_path or not os.path.exists(result_path):
        return "File not found", 404
    
    # Log conversion before download
    log_conversion('pdf-to-ppt', result_filename, 'presentation.pptx', 'file', result_path, 'success')
    
    # Upload to Cloudinary
    username = session.get('username')
    cloudinary_folder = f'storage/{username}/files'
    
    try:
        cloudinary_result = cloudinary.uploader.upload(result_path, 
                                                      folder=cloudinary_folder,
                                                      resource_type='raw',
                                                      public_id=f'ppt_{uuid.uuid4().hex}.pptx')
        cloudinary_url = cloudinary_result['secure_url']
        store_url_in_firebase(cloudinary_url, 'files', 'presentation.pptx')
    except Exception as e:
        print(f"Cloudinary upload error: {e}")
    
    # Background cleanup
    def cleanup():
        import time
        time.sleep(2)
        if os.path.exists(result_path):
            os.remove(result_path)
        session.pop('ppt_result_id', None)
        session.pop('ppt_result_path', None)
        session.pop('ppt_result_filename', None)
    
    cleanup_thread = threading.Thread(target=cleanup, daemon=True)
    cleanup_thread.start()
    
    base_name = result_filename.rsplit('.', 1)[0] if '.' in result_filename else result_filename
    download_name = f"{base_name}.pptx"
    
    return send_file(result_path, as_attachment=True, download_name=download_name,
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

# --- Background Remover ---
@app.route('/remove_background', methods=['POST'])
@login_required
def remove_background():
    logger.debug("Received request at /remove_background")
 
    if 'image' not in request.files:
        logger.error("No image uploaded")
        return jsonify({'error': 'No image uploaded!'}), 400
 
    file = request.files['image']
    valid_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp')
 
    if not file.filename.lower().endswith(valid_extensions):
        return jsonify({'error': 'Unsupported file type. Please upload a JPG, PNG, GIF, BMP, TIFF, or WebP image.'}), 400
 
    try:
        # ── 1. Read raw bytes (reliable size check) ──────────────────
        file.stream.seek(0)
        image_bytes = file.stream.read()
        if len(image_bytes) == 0:
            return jsonify({'error': 'Uploaded file is empty.'}), 400
        if len(image_bytes) > 10 * 1024 * 1024:
            return jsonify({'error': 'File size exceeds 10 MB limit.'}), 400
 
        # ── 2. Validate image can be opened ──────────────────────────
        try:
            image = Image.open(io.BytesIO(image_bytes))
            image.verify()                          # catch truncated / corrupt files
            image = Image.open(io.BytesIO(image_bytes))   # re-open after verify()
        except Exception as img_err:
            logger.error(f"Image validation failed: {img_err}")
            return jsonify({'error': f'Invalid or corrupted image file: {img_err}'}), 400
 
        # ── 3. Optional resize for speed ─────────────────────────────
        MAX_DIM = 1024
        if image.size[0] > MAX_DIM or image.size[1] > MAX_DIM:
            image.thumbnail((MAX_DIM, MAX_DIM), Image.LANCZOS)
            logger.info(f"Image resized to {image.size} for faster bg removal")
 
        if image.mode != 'RGBA':
            image = image.convert('RGBA')
 
        # ── 4. rembg — remove background ────────────────────────────
        try:
            from rembg import remove as rembg_remove
        except ImportError:
            logger.error("rembg library is not installed")
            return jsonify({'error': 'Background removal library (rembg) is not installed on the server.'}), 500
 
        try:
            _buf_in = io.BytesIO()
            image.save(_buf_in, format='PNG')
            _buf_in.seek(0)
            output_bytes = rembg_remove(_buf_in.getvalue())
        except Exception as rembg_err:
            logger.error(f"rembg processing failed: {rembg_err}", exc_info=True)
            return jsonify({'error': f'Background removal processing failed: {rembg_err}'}), 500
 
        # ── 5. Load result image ─────────────────────────────────────
        try:
            output_image = Image.open(io.BytesIO(output_bytes)).convert('RGBA')
        except Exception as load_err:
            logger.error(f"Failed to load rembg output: {load_err}")
            return jsonify({'error': 'Background removal produced invalid output.'}), 500
 
        # ── 6. Save result to temp file for download endpoint ─────────
        output_filename = f"bg_removed_{uuid.uuid4().hex}.png"
        OUTPUT_FOLDER = 'converted'
        os.makedirs(OUTPUT_FOLDER, exist_ok=True)
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        output_image.save(output_path, format='PNG')
        logger.debug(f"Result saved to: {output_path}")
 
        # ── 7. Build base64 preview for instant display in browser ────
        preview_buf = io.BytesIO()
        output_image.save(preview_buf, format='PNG')
        preview_buf.seek(0)
        preview_b64 = base64.b64encode(preview_buf.getvalue()).decode('utf-8')
        preview_data_url = f"data:image/png;base64,{preview_b64}"
 
        # ── 8. Upload to Cloudinary ───────────────────────────────────
        cloudinary_url = None
        try:
            img_buffer = io.BytesIO()
            output_image.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            username = session.get('username')
            cloudinary_folder = f'storage/{username}/img'
            cloudinary_result = cloudinary.uploader.upload(
                img_buffer,
                folder=cloudinary_folder,
                resource_type='raw',
                public_id=output_filename
            )
            cloudinary_url = cloudinary_result['secure_url']
            store_url_in_firebase(cloudinary_url, 'img', output_filename)
        except Exception as cloud_err:
            logger.warning(f"Cloudinary upload failed (non-fatal): {cloud_err}")
 
        log_conversion('background-remover', file.filename, output_filename, output_path, cloudinary_url)
 
        # ── 9. Schedule cleanup of the local file (5 minutes) ─────────
        def _cleanup_file(path, delay=300):
            time.sleep(delay)
            try:
                if os.path.exists(path):
                    os.remove(path)
                    logger.debug(f"Cleaned up temp file: {path}")
            except Exception as e:
                logger.warning(f"Cleanup failed for {path}: {e}")
 
        threading.Thread(target=_cleanup_file, args=(output_path,), daemon=True).start()
 
        logger.debug("Background removal successful")
        return jsonify({
            'success': True,
            'download_url': f'/download_bg_removed/{output_filename}',
            'preview_url': preview_data_url,          # ← inline base64 for instant display
            'filename': output_filename,
            'cloudinary_url': cloudinary_url,
            'message': 'Background removed successfully!'
        }), 200
 
    except Exception as e:
        logger.error(f"Background removal failed: {e}", exc_info=True)
        return jsonify({'error': f'Background removal failed: {e}'}), 500
 
 
@app.route('/download_bg_removed/<filename>', methods=['GET'])
@login_required
def download_bg_removed(filename):
    try:
        # Prevent path traversal
        if '..' in filename or '/' in filename or '\\' in filename:
            logger.error(f"Invalid filename attempt: {filename}")
            return jsonify({'error': 'Invalid filename'}), 400
 
        OUTPUT_FOLDER = 'converted'
        file_path = os.path.join(OUTPUT_FOLDER, filename)
 
        if not os.path.exists(file_path):
            # Try Cloudinary as fallback
            logger.warning(f"Local file not found, attempting Cloudinary fallback: {file_path}")
            username = session.get('username')
            safe_key = re.sub(r'[./#$\[\]]', '_', filename)
            ref = db.reference(f'storage/{username}/img/{safe_key}')
            data = ref.get()
            if data and 'url' in data:
                r = requests.get(data['url'], timeout=15)
                if r.status_code == 200:
                    return send_file(
                        io.BytesIO(r.content),
                        as_attachment=True,
                        download_name='background_removed.png',
                        mimetype='image/png'
                    )
            return jsonify({'error': 'File not found. It may have expired — please process the image again.'}), 404
 
        logger.debug(f"Downloading file: {file_path}")
        return send_file(
            file_path,
            as_attachment=True,
            download_name='background_removed.png',
            mimetype='image/png'
        )
 
    except Exception as e:
        logger.error(f"Download failed: {e}", exc_info=True)
        return jsonify({'error': f'Download failed: {e}'}), 500
 

# --- Plagiarism Scanner ---

def fetch_web_snippets(query, max_results=5):
    """Scrape DuckDuckGo for snippets."""
    search_url = f"https://html.duckduckgo.com/html/?q={query}"
    headers = {"User-Agent": "Mozilla/5.0"}
    try:
        response = requests.get(search_url, headers=headers, timeout=5)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        results = soup.find_all('a', {'class': 'result__a'}, limit=max_results)
        snippets = [r.get_text(strip=True) for r in results if r.get_text(strip=True)]
        return snippets
    except Exception as e:
        logger.error(f"Error fetching web snippets: {str(e)}")
        return []

def call_openrouter_similarity(text_a, text_b):
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
        "User-Agent": "Mozilla/5.0"
    }
    messages = [
        {
            "role": "system",
            "content": ("You are a plagiarism detection assistant. Given two texts, "
                        "respond ONLY with a plagiarism similarity percentage (0 to 100) and a brief explanation, separated by a newline.")
        },
        {
            "role": "user",
            "content": f"Text A:\n{text_a}\n\nText B:\n{text_b}"
        }
    ]
    data = {
        "model": OPENROUTER_MODEL,
        "messages": messages
    }
    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data)
        response.raise_for_status()
        return response.json()['choices'][0]['message']['content'].strip()
    except Exception as e:
        logger.error(f"OpenRouter API error in similarity check: {str(e)}")
        return f"Error: {str(e)}"

@app.route('/check_plagiarism', methods=['POST'])
@login_required
def check_plagiarism():
    try:
        input_text = request.form.get('text', '').strip()
        if not input_text:
            return jsonify({'error': 'No text provided'}), 400
        if len(input_text) < 20:
            results = [{"snippet": "", "similarity": "Input text too short to check plagiarism."}]
            return jsonify({'results': results})

        query = input_text[:100]
        snippets = fetch_web_snippets(query)
        results = []
        for snippet in snippets:
            sim = call_openrouter_similarity(input_text, snippet)
            results.append({"snippet": snippet, "similarity": sim})
        return jsonify({'results': results})
    except Exception as e:
        logger.error(f"Plagiarism check failed: {str(e)}")
        return jsonify({'error': f"Plagiarism check failed: {str(e)}"}), 500

# --- Text to Speech ---

@app.route('/generate_tts', methods=['POST'])
@login_required
def generate_tts():
    data = request.get_json()
    text = data.get('text', '').strip()
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    try:
        filename = f"{uuid.uuid4()}.mp3"
        
        # Use temporary file instead of audio folder
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as temp_file:
            filepath = temp_file.name
            
        tts = gTTS(text=text, lang='en')
        tts.save(filepath)

        # Upload to cloudinary first, before logging
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/audio'
        cloudinary_url = upload_to_cloudinary(filepath, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'audio', filename)

        # Now log the conversion with the cloudinary_url
        log_conversion('text-to-speech', 'user_input.txt', filename, filepath, cloudinary_url)

        # Store the audio data for download in cache
        with open(filepath, 'rb') as f:
            audio_data = f.read()
        
        # Store in audio cache for immediate download
        audio_cache[filename] = audio_data
        
        # Clean up temp file
        try:
            os.remove(filepath)
        except Exception as e:
            logger.warning(f"Failed to clean up temp audio file: {str(e)}")

        return jsonify({
            'success': True,
            'audio_url': f'/stream_audio/{filename}',  # For playing in browser
            'download_url': f'/download_audio/{filename}',  # For downloading
            'filename': filename,
            'message': 'Text converted to speech successfully!'
        })
    except Exception as e:
        logger.error(f"Text to speech conversion failed: {str(e)}", exc_info=True)
        return jsonify({'error': f"Text to speech conversion failed: {str(e)}"}), 500

@app.route('/download_audio/<filename>')
@login_required
def download_audio(filename):
    logger.info(f"Audio download requested for: {filename}")
    logger.info(f"Available audio files in cache: {list(audio_cache.keys())}")
    
    # Check if file exists in memory cache
    if filename not in audio_cache:
        logger.error(f"Audio file {filename} not found in cache")
        return jsonify({'error': 'Audio file not found or expired'}), 404
    
    try:
        # Get audio data from memory cache
        audio_data = audio_cache[filename]
        logger.info(f"Found audio data, size: {len(audio_data)} bytes")
        
        # Create a BytesIO object with the audio data
        audio_buffer = io.BytesIO(audio_data)
        audio_buffer.seek(0)
        
        # Don't delete from cache immediately - let it be accessed multiple times
        # The cache will be cleaned up by a timer or manually
        
        logger.info(f"Serving audio file: {filename}")
        return send_file(
            audio_buffer, 
            as_attachment=True, 
            download_name=filename, 
            mimetype='audio/mpeg'
        )
    except Exception as e:
        logger.error(f"Audio download failed: {str(e)}", exc_info=True)
        return jsonify({'error': f"Audio download failed: {str(e)}"}), 500

@app.route('/stream_audio/<filename>')
@login_required
def stream_audio(filename):
    """Stream audio for playing in browser (not as download)"""
    logger.info(f"Audio stream requested for: {filename}")
    
    # Check if file exists in memory cache
    if filename not in audio_cache:
        logger.error(f"Audio file {filename} not found in cache")
        return jsonify({'error': 'Audio file not found or expired'}), 404
    
    try:
        # Get audio data from memory cache
        audio_data = audio_cache[filename]
        logger.info(f"Streaming audio data, size: {len(audio_data)} bytes")
        
        # Create a BytesIO object with the audio data
        audio_buffer = io.BytesIO(audio_data)
        audio_buffer.seek(0)
        
        logger.info(f"Streaming audio file: {filename}")
        return send_file(
            audio_buffer, 
            as_attachment=False,  # Don't force download for streaming
            download_name=filename, 
            mimetype='audio/mpeg'
        )
    except Exception as e:
        logger.error(f"Audio streaming failed: {str(e)}", exc_info=True)
        return jsonify({'error': f"Audio streaming failed: {str(e)}"}), 500

# --- Speech to Text ---

@app.route('/save_transcript', methods=['POST'])
@login_required
def save_transcript():
    try:
        data = request.get_json()
        transcript = data.get('transcript', '')
        if not transcript:
            return jsonify({'error': 'No transcript provided'}), 400
        transcript_filename = f"transcript_{uuid.uuid4().hex}.txt"
        
        # Use temporary file for transcript
        with tempfile.NamedTemporaryFile(delete=False, suffix='.txt', mode='w', encoding='utf-8') as temp_file:
            temp_file.write(transcript + '\n')
            transcript_path = temp_file.name
            
        # Upload to cloudinary first, before logging
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/txt'
        cloudinary_url = upload_to_cloudinary(transcript_path, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'txt', transcript_filename)

        # Now log the conversion with the cloudinary_url
        log_conversion('speech-to-text', 'transcript.txt', transcript_filename, transcript_path, cloudinary_url)

        # Clean up temp file
        try:
            os.remove(transcript_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp transcript file: {str(e)}")

        return jsonify({'message': 'Transcript saved successfully'})
    except Exception as e:
        return jsonify({'error': f"Transcript save failed: {str(e)}"}), 500

@app.route('/upload_audio', methods=['POST'])
@login_required
def upload_audio():
    temp_file_path = None
    wav_path = None
    try:
        if 'audioFile' not in request.files:
            return jsonify({'error': 'No audio file provided'}), 400
        
        audio_file = request.files['audioFile']
        if audio_file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
            
        # Create temporary file for the uploaded audio
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as temp_file:
            audio_file.save(temp_file)
            temp_file_path = temp_file.name
        
        # Try to convert using pydub (requires FFmpeg)
        try:
            audio = AudioSegment.from_mp3(temp_file_path)
            wav_path = temp_file_path.replace('.mp3', '.wav')
            audio.export(wav_path, format='wav')
        except Exception as convert_error:
            # If pydub fails, try direct wav loading if file is already wav format
            if audio_file.filename.lower().endswith('.wav'):
                wav_path = temp_file_path
            else:
                # Fallback: try to use the mp3 file directly with speech_recognition
                try:
                    with sr.AudioFile(temp_file_path) as source:
                        audio_data = recognizer.record(source)
                        transcript = recognizer.recognize_google(audio_data)
                        
                    # Cleanup and return
                    if temp_file_path and os.path.exists(temp_file_path):
                        os.remove(temp_file_path)
                    log_conversion('speech-to-text', audio_file.filename, 'transcript.txt', None, None, 'success')
                    return jsonify({'transcript': transcript})
                except Exception:
                    return jsonify({'error': 'FFmpeg is required for MP3 conversion. Please install FFmpeg or upload a WAV file instead.'}), 400
        
        # Process the wav file for speech recognition
        with sr.AudioFile(wav_path) as source:
            # Adjust for ambient noise
            recognizer.adjust_for_ambient_noise(source, duration=0.5)
            audio_data = recognizer.record(source)
            
            try:
                transcript = recognizer.recognize_google(audio_data)
                if not transcript.strip():
                    transcript = "No speech detected in the audio file"
            except sr.UnknownValueError:
                transcript = "Could not understand the audio. Please ensure the audio is clear and contains speech."
            except sr.RequestError as e:
                transcript = f"Speech recognition service error: {str(e)}"
        
        log_conversion('speech-to-text', audio_file.filename, 'transcript.txt', None, None, 'success')
        return jsonify({'transcript': transcript})
        
    except Exception as e:
        logger.error(f"Speech to text conversion error: {str(e)}")
        return jsonify({'error': f"Speech to text conversion failed: {str(e)}"}), 500
    finally:
        # Cleanup temporary files
        try:
            if temp_file_path and os.path.exists(temp_file_path):
                os.remove(temp_file_path)
            if wav_path and wav_path != temp_file_path and os.path.exists(wav_path):
                os.remove(wav_path)
        except Exception as cleanup_error:
            logger.warning(f"Failed to cleanup temp files: {str(cleanup_error)}")

# --- AI PDF Editor ---

def extract_structured_text(path):
    try:
        doc = fitz.open(path)
        structured_data = []
        for page in doc:
            blocks = page.get_text("dict")['blocks']
            for block in blocks:
                if 'lines' in block:
                    for line in block['lines']:
                        line_text = " ".join([span['text'] for span in line['spans']])
                        structured_data.append(line_text)
        doc.close()
        return "\n".join(structured_data)
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        return f"Error extracting text from PDF: {str(e)}"

def retry_api_request(url, headers, data, max_retries=3, delay=2):
    for attempt in range(max_retries):
        try:
            response = requests.post(url, headers=headers, json=data, timeout=30)
            if response.status_code == 429:
                logger.warning(f"Rate limit hit on attempt {attempt + 1}, retrying after {delay} seconds")
                time.sleep(delay)
                delay *= 2
                continue
            response.raise_for_status()
            return response
        except requests.exceptions.RequestException as e:
            logger.error(f"API request failed on attempt {attempt + 1}: {str(e)}")
            if attempt == max_retries - 1:
                raise e
            time.sleep(delay)
            delay *= 2
    raise Exception("Max retries exceeded for API request")

@app.route('/analyze', methods=['POST'])
@login_required
def analyze():
    global latest_text
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    uploaded_file = request.files['file']
    if not uploaded_file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Only PDF files are supported'}), 400
    filename = secure_filename(uploaded_file.filename)
    
    # Use temporary file instead of upload folder
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
        uploaded_file.save(temp_file)
        file_path = temp_file.name
    try:
        extracted_text = extract_structured_text(file_path)
        if extracted_text.startswith('Error'):
            os.remove(file_path)
            return jsonify({'error': extracted_text}), 500
        if not extracted_text.strip():
            os.remove(file_path)
            return jsonify({'error': 'No text extracted from PDF. Ensure the PDF contains selectable text, not images.'}), 400
        latest_text = extracted_text
        prompt = (
            "You are an intelligent assistant analyzing a PDF document. Your task is to identify blank or unfilled fields such as "
            "'Date: ____', 'Name: ________', 'Signature: [____]', or other placeholders (e.g., '________', '[____]', empty lines after labels). "
            "For each identified field, suggest a reasonable completion based on context (e.g., use today's date 'June 25, 2025' for date fields, 'John Doe' for name fields, 'Signature' for signature fields). "
            "Return the results in the following format:\n"
            "Identified Fields:\n"
            "- Field: [Description], Suggestion: [Suggested Value]\n"
            "If no blank fields are found, state: 'No blank or unfilled fields detected.'\n\n"
            f"Document Text:\n{extracted_text[:2000]}"
        )

        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }

        data = {
            "model": OPENROUTER_MODEL,
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 1000
        }

        response = retry_api_request(OPENROUTER_API_URL, headers, data)
        result = response.json()
        suggestions = result.get('choices', [{}])[0].get('message', {}).get('content', '')

        if not suggestions:
            raise ValueError("Empty response content from API")

        log_conversion('ai-pdf-editor', filename, 'analysis.json', None, None, 'success')
        return jsonify({"text": extracted_text, "suggestions": suggestions})
    except requests.exceptions.RequestException as e:
        logger.error(f"OpenRouter API request failed: {str(e)}")
        return jsonify({'error': f"Failed to analyze document: API request error - {str(e)}"}), 500
    except ValueError as e:
        logger.error(f"OpenRouter API response error: {str(e)}")
        return jsonify({'error': f"Failed to analyze document: Invalid API response - {str(e)}"}), 500
    except Exception as e:
        logger.error(f"PDF analysis failed: {str(e)}")
        return jsonify({'error': f"Failed to analyze document: {str(e)}"}), 500
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)

@app.route('/edit', methods=['POST'])
@login_required
def edit():
    global latest_text
    try:
        updated_text = request.json.get('updated_text')
        if not updated_text:
            return jsonify({'error': 'No updated text provided'}), 400
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        for line in updated_text.split('\n'):
            pdf.multi_cell(0, 10, line)
        output_filename = f"edited_{uuid.uuid4().hex}.pdf"
        
        # Use temporary file instead of upload folder
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            output_path = temp_file.name
            
        pdf.output(output_path)
        log_conversion('ai-pdf-editor', 'user_input.txt', output_filename, output_path, None, 'success')
        return send_file(
            output_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name="edited_document.pdf"
        )
    except Exception as e:
        logger.error(f"PDF edit failed: {str(e)}")
        return jsonify({'error': f"Failed to edit document: {str(e)}"}), 500
    finally:
        # Clean up temp file
        try:
            if 'output_path' in locals() and os.path.exists(output_path):
                os.remove(output_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp edit file: {str(e)}")

@app.route('/fill_from_prompt', methods=['POST'])
@login_required
def fill_from_prompt():
    global latest_text
    try:
        user_prompt = request.json.get('user_prompt')
        current_text = request.json.get('current_text', '')
        
        if not user_prompt:
            return jsonify({'error': 'No user prompt provided'}), 400
        
        # Use current_text if provided, otherwise fall back to latest_text
        text_to_process = current_text if current_text.strip() else latest_text
        
        if not text_to_process:
            return jsonify({'error': 'No text available to process. Please analyze a document first or add some text.'}), 400
            
        ai_instruction = (
            "Based on the following document text, a user wants to add or update content as follows:\n"
            f"Instruction: {user_prompt}\n\n"
            f"Document:\n{text_to_process[:2000]}\n\n"
            "Provide the updated version of the document with the user request applied appropriately. "
            "Return only the revised document content."
        )
        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }
        data = {
            "model": OPENROUTER_MODEL,
            "messages": [{"role": "user", "content": ai_instruction}],
            "max_tokens": 1000
        }
        response = retry_api_request(OPENROUTER_API_URL, headers, data)
        result = response.json()
        updated_text = result.get('choices', [{}])[0].get('message', {}).get('content', '')
        if not updated_text:
            return jsonify({'error': 'Empty response from API'}), 500
        latest_text = updated_text
        log_conversion('ai-pdf-editor', 'user_prompt.txt', 'updated_text.txt', None, None, 'success')
        return jsonify({"updated_text": updated_text})
    except Exception as e:
        logger.error(f"Prompt-based edit failed: {str(e)}")
        return jsonify({'error': f"Prompt-based edit failed: {str(e)}"}), 500

# --- Text Summarizer ---

@app.route('/summarize', methods=['POST'])
@login_required
def summarize():
    try:
        data = request.get_json()
        text = data.get('text', '').strip()
        if not text:
            logger.error("No text provided to summarizer endpoint")
            return jsonify({'error': 'No text provided'}), 400
        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": OPENROUTER_MODEL,
            "messages": [
                {
                    "role": "user",
                    "content": (
                        "Summarize the following text in 3 sentences. "
                        "Only return the summary content. Do not add any introduction, title, or prefix:\n\n"
                        f"{text[:2000]}"
                    )
                }
            ],
            "max_tokens": 500
        }
        response = retry_api_request(OPENROUTER_API_URL, headers, payload)
        result = response.json()
        summary = result.get('choices', [{}])[0].get('message', {}).get('content', '').strip()
        if not summary:
            logger.error(f"Empty summary from API. Response: {json.dumps(result)}")
            return jsonify({'error': 'Empty summary from API'}), 500

        summary_filename = f"summary_{uuid.uuid4().hex}.txt"
        
        # Use temporary file for summary
        with tempfile.NamedTemporaryFile(delete=False, suffix='.txt', mode='w', encoding='utf-8') as temp_file:
            temp_file.write(summary)
            summary_path = temp_file.name

        # Upload to cloudinary first, before logging
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/txt'
        try:
            cloudinary_url = upload_to_cloudinary(summary_path, cloudinary_folder)
            store_url_in_firebase(cloudinary_url, 'txt', summary_filename)
        except Exception as upload_e:
            logger.error(f"Cloudinary/Firebase upload failed for summary: {str(upload_e)}")
            cloudinary_url = None  # Set to None if upload fails

        # Now log the conversion with the cloudinary_url
        log_conversion('text-summarizer', 'user_input.txt', summary_filename, summary_path, cloudinary_url)

        # Clean up temp file
        try:
            os.remove(summary_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp summary file: {str(e)}")

        return jsonify({'summary': summary})
    except Exception as e:
        logger.error(f"Text summarization failed: {str(e)}", exc_info=True)
        return jsonify({'error': f"Text summarization failed: {str(e)}"}), 500

@app.route('/admin/logs', methods=['GET'])
@login_required
def get_conversion_logs():
    """API endpoint to fetch conversion logs for the history page."""
    try:
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        c.execute('''
            SELECT id, conversion_type, original_filename, converted_filename, file_path, timestamp, cloudinary_url, username, status
            FROM conversions 
            WHERE username = ?
            ORDER BY timestamp DESC
        ''', (session.get('username'),))
        rows = c.fetchall()
        conn.close()
        
        logs = []
        for row in rows:
            # Use the actual status from database, with fallback logic
            actual_status = row[8] if len(row) > 8 and row[8] else None
            if actual_status is None:
                # Fallback: determine status based on available data
                if row[6]:  # cloudinary_url exists
                    actual_status = 'success'
                elif row[1] in ['speech-to-text', 'document-screener', 'ai-pdf-editor', 'text-summarizer']:
                    # These don't always need cloudinary_url to be successful
                    actual_status = 'success'
                else:
                    actual_status = 'error'
            
            log_entry = {
                'id': row[0],
                'conversion_type': row[1],
                'original_filename': row[2],
                'filename': row[3],
                'download_path': f"/download-file/{row[0]}" if (row[6] or actual_status == 'success') else None,
                'timestamp': row[5],
                'status': actual_status,
                'file_size': 'Unknown',  # You can add file size calculation if needed
                'cloudinary_url': row[6],
                'username': row[7]
            }
            logs.append(log_entry)
        
        return jsonify({'success': True, 'logs': logs})
    except Exception as e:
        logger.error(f"Failed to fetch conversion logs: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500

# --- Cache Cleanup Functions ---

def cleanup_audio_cache():
    """Clean up audio cache to prevent memory leaks"""
    try:
        if audio_cache:
            logger.info(f"Cleaning up {len(audio_cache)} audio files from cache")
            audio_cache.clear()
    except Exception as e:
        logger.error(f"Error cleaning up audio cache: {str(e)}")

# Schedule periodic cleanup every 10 minutes
import threading
import time

def periodic_cleanup():
    while True:
        time.sleep(600)  # 10 minutes
        cleanup_audio_cache()

# Start cleanup thread
cleanup_thread = threading.Thread(target=periodic_cleanup, daemon=True)
cleanup_thread.start()

# --- Test Routes for Debugging ---

@app.route('/test_dependencies')
def test_dependencies():
    """Test if all required dependencies are working"""
    try:
        # Test PIL/Pillow
        from PIL import Image
        test_image = Image.new('RGB', (100, 100), color='red')
        
        # Test Cloudinary
        import cloudinary
        cloudinary_status = "Configured" if cloudinary.config().cloud_name else "Not configured"
        
        # Test Firebase
        firebase_status = "Connected" if db else "Not connected"
        
        return jsonify({
            'success': True,
            'dependencies': {
                'PIL': 'Working',
                'Cloudinary': cloudinary_status,
                'Firebase': firebase_status
            }
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        })

# --- Profile Management Routes ---

def update_users_table():
    """Update users table to include profile fields"""
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    
    # Add new columns if they don't exist
    profile_columns = [
        ('name', 'TEXT'),
        ('email', 'TEXT'),
        ('contact_number', 'TEXT'),
        ('country', 'TEXT'),
        ('profile_picture', 'TEXT'),
        ('membership_status', 'TEXT DEFAULT "Standard"'),
        ('created_at', 'TIMESTAMP DEFAULT CURRENT_TIMESTAMP')
    ]
    
    for column_name, column_type in profile_columns:
        try:
            c.execute(f'ALTER TABLE users ADD COLUMN {column_name} {column_type}')
        except sqlite3.OperationalError:
            pass  # Column already exists
    
    conn.commit()
    conn.close()

def get_user_profile(username):
    """Get user profile data from database and Firebase"""
    try:
        # Get from Firebase first
        user_ref = db.reference(f'Data/{username}')
        user_data = user_ref.get()
        
        if user_data:
            return user_data
        
        # Fallback to SQLite
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        c.execute('''SELECT username, name, email, contact_number, country, 
                           profile_picture, membership_status, created_at 
                    FROM users WHERE username = ?''', (username,))
        row = c.fetchone()
        conn.close()
        
        if row:
            return {
                'username': row[0],
                'name': row[1],
                'email': row[2],
                'contact_number': row[3],
                'country': row[4],
                'profile_picture': row[5],
                'membership_status': row[6] or 'Standard',
                'created_at': row[7]
            }
        return None
    except Exception as e:
        logger.error(f"Error getting user profile: {str(e)}")
        return None

def get_user_context():
    """Get user context for templates"""
    username = session.get('username')
    user = get_user_profile(username)
    
    if not user:
        user = {'username': username, 'profile_picture': None}
    
    return {'user': user}

def update_user_profile(username, profile_data):
    """Update user profile in both Firebase and SQLite"""
    try:
        # Update Firebase
        user_ref = db.reference(f'Data/{username}')
        existing_data = user_ref.get() or {}
        existing_data.update(profile_data)
        user_ref.set(existing_data)
        
        # Update SQLite as backup
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        
        # Check if user exists in SQLite
        c.execute('SELECT id FROM users WHERE username = ?', (username,))
        user_exists = c.fetchone()
        
        if user_exists:
            # Update existing user
            c.execute('''UPDATE users SET name = ?, email = ?, contact_number = ?, 
                               country = ?, profile_picture = ?, membership_status = ?
                        WHERE username = ?''', 
                     (profile_data.get('name'), profile_data.get('email'), 
                      profile_data.get('contact_number'), profile_data.get('country'),
                      profile_data.get('profile_picture'), profile_data.get('membership_status'),
                      username))
        else:
            # Insert new user record
            c.execute('''INSERT INTO users (username, name, email, contact_number, 
                               country, profile_picture, membership_status)
                        VALUES (?, ?, ?, ?, ?, ?, ?)''',
                     (username, profile_data.get('name'), profile_data.get('email'),
                      profile_data.get('contact_number'), profile_data.get('country'),
                      profile_data.get('profile_picture'), profile_data.get('membership_status')))
        
        conn.commit()
        conn.close()
        return True
    except Exception as e:
        logger.error(f"Error updating user profile: {str(e)}")
        return False

@app.route('/profile')
@login_required
def profile():
    """Display user profile page"""
    username = session.get('username')
    user = get_user_profile(username)
    
    if not user:
        user = {'username': username}
    
    return render_template('profile.html', user=user)

@app.route('/update_profile', methods=['POST'])
@login_required
def update_profile():
    """Update user profile information"""
    try:
        username = session.get('username')
        
        # Get form data
        profile_data = {
            'name': request.form.get('name', '').strip(),
            'username': request.form.get('username', '').strip(),
            'email': request.form.get('email', '').strip(),
            'contact_number': request.form.get('contact_number', '').strip(),
            'country': request.form.get('country', '').strip(),
            'membership_status': request.form.get('membership_status', 'Standard')
        }
        
        # Validate required fields
        required_fields = ['name', 'username', 'email', 'country']
        for field in required_fields:
            if not profile_data[field]:
                return render_template('profile.html', 
                                     user=get_user_profile(username),
                                     error_message=f'{field.replace("_", " ").title()} is required')
        
        # Validate email format
        import re
        email_pattern = r'^[^\s@]+@[^\s@]+\.[^\s@]+$'
        if not re.match(email_pattern, profile_data['email']):
            return render_template('profile.html', 
                                 user=get_user_profile(username),
                                 error_message='Please enter a valid email address')
        
        # Check if username is being changed and if new username exists
        if profile_data['username'] != username:
            existing_user = get_user_profile(profile_data['username'])
            if existing_user:
                return render_template('profile.html', 
                                     user=get_user_profile(username),
                                     error_message='Username already exists')
        
        # Update profile
        if update_user_profile(username, profile_data):
            # Update session if username changed
            if profile_data['username'] != username:
                session['username'] = profile_data['username']
            
            return render_template('profile.html', 
                                 user=get_user_profile(profile_data['username']),
                                 success_message='Profile updated successfully!')
        else:
            return render_template('profile.html', 
                                 user=get_user_profile(username),
                                 error_message='Failed to update profile')
    
    except Exception as e:
        logger.error(f"Profile update error: {str(e)}")
        return render_template('profile.html', 
                             user=get_user_profile(session.get('username')),
                             error_message='An error occurred while updating profile')

@app.route('/upload_profile_picture', methods=['POST'])
@login_required
def upload_profile_picture():
    """Handle profile picture upload"""
    try:
        username = session.get('username')
        logger.info(f"Profile picture upload request from user: {username}")
        
        if 'profile_picture' not in request.files:
            logger.error("No profile_picture in request.files")
            return jsonify({'success': False, 'error': 'No file uploaded'})
        
        file = request.files['profile_picture']
        logger.info(f"File received: {file.filename}, size: {file.content_length}")
        
        if file.filename == '':
            logger.error("Empty filename")
            return jsonify({'success': False, 'error': 'No file selected'})
        
        # Validate file type
        allowed_extensions = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
        file_extension = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        
        if file_extension not in allowed_extensions:
            logger.error(f"Invalid file extension: {file_extension}")
            return jsonify({'success': False, 'error': 'Invalid file type. Please upload an image file.'})
        
        # Create temporary file for processing
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
            file.save(temp_file.name)
            temp_path = temp_file.name
            logger.info(f"File saved to temp path: {temp_path}")
        
        try:
            # Check if PIL/Pillow is working
            logger.info("Processing image with PIL...")
            
            # Resize and optimize image
            with Image.open(temp_path) as img:
                logger.info(f"Image opened: {img.size}, mode: {img.mode}")
                
                # Convert to RGB if necessary
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                    logger.info("Image converted to RGB")
                
                # Resize to max 500x500 while maintaining aspect ratio
                img.thumbnail((500, 500), Image.Resampling.LANCZOS)
                logger.info(f"Image resized to: {img.size}")
                
                # Save optimized image
                optimized_path = temp_path.replace(f'.{file_extension}', '_optimized.jpg')
                img.save(optimized_path, 'JPEG', quality=85, optimize=True)
                logger.info(f"Optimized image saved to: {optimized_path}")
            
            # Upload to Cloudinary
            logger.info("Uploading to Cloudinary...")
            cloudinary_folder = f'storage/{username}/profile'
            
            result = cloudinary.uploader.upload(
                optimized_path,
                folder=cloudinary_folder,
                public_id=f'profile_picture_{username}',
                overwrite=True,
                resource_type='image',
                format='jpg'
            )
            
            profile_picture_url = result['secure_url']
            logger.info(f"Cloudinary upload successful: {profile_picture_url}")
            
            # Update user profile with new picture URL
            profile_data = {'profile_picture': profile_picture_url}
            if update_user_profile(username, profile_data):
                logger.info("Profile updated successfully")
                return jsonify({
                    'success': True, 
                    'profile_picture_url': profile_picture_url,
                    'message': 'Profile picture updated successfully!'
                })
            else:
                logger.error("Failed to update user profile in database")
                return jsonify({'success': False, 'error': 'Failed to save profile picture URL'})
        
        except Exception as upload_error:
            logger.error(f"Image upload error: {str(upload_error)}", exc_info=True)
            return jsonify({'success': False, 'error': f'Failed to process and upload image: {str(upload_error)}'})
        
        finally:
            # Clean up temporary files
            try:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
                    logger.info(f"Cleaned up temp file: {temp_path}")
                if 'optimized_path' in locals() and os.path.exists(optimized_path):
                    os.remove(optimized_path)
                    logger.info(f"Cleaned up optimized file: {optimized_path}")
            except Exception as cleanup_error:
                logger.warning(f"Failed to clean up temp files: {str(cleanup_error)}")
    
    except Exception as e:
        logger.error(f"Profile picture upload error: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'error': f'An error occurred while uploading the image: {str(e)}'})

@app.route('/remove_profile_picture', methods=['POST'])
@login_required
def remove_profile_picture():
    try:
        username = session.get('username')
        
        # Update profile picture to null in user profile
        profile_data = {'profile_picture': None}
        if update_user_profile(username, profile_data):
            return jsonify({'success': True, 'message': 'Profile picture removed successfully!'})
        else:
            return jsonify({'success': False, 'error': 'Failed to remove profile picture'})
    except Exception as e:
        logger.error(f"Remove profile picture error: {str(e)}")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/change_password')
@login_required
def change_password():
    """Display change password form"""
    return render_template('change_password.html')

@app.route('/update_password', methods=['POST'])
@login_required
def update_password():
    """Update user password"""
    try:
        username = session.get('username')
        current_password = request.form.get('current_password')
        new_password = request.form.get('new_password')
        confirm_password = request.form.get('confirm_password')
        
        # Validate inputs
        if not all([current_password, new_password, confirm_password]):
            return render_template('change_password.html', 
                                 error='All fields are required')
        
        if new_password != confirm_password:
            return render_template('change_password.html', 
                                 error='New passwords do not match')
        
        if len(new_password) < 6:
            return render_template('change_password.html', 
                                 error='Password must be at least 6 characters long')
        
        # Verify current password
        user_data = get_user_by_username(username)
        if not user_data or not check_password_hash(user_data['password'], current_password):
            return render_template('change_password.html', 
                                 error='Current password is incorrect')
        
        # Update password in Firebase
        new_password_hash = generate_password_hash(new_password)
        cred_ref = db.reference(f'credentials/users/{username}')
        cred_ref.update({'password': new_password_hash})
        
        # Also update in user data
        user_ref = db.reference(f'Data/{username}')
        user_ref.update({'password': new_password_hash})
        
        return render_template('change_password.html', 
                             success='Password updated successfully!')
    
    except Exception as e:
        logger.error(f"Password update error: {str(e)}")
        return render_template('change_password.html', 
                             error='An error occurred while updating password')

@app.route('/email-settings')
@login_required
def email_settings():
    """Display email settings page for phone/email verification"""
    username = session.get('username')
    user_data = get_user_profile(username)
    
    if not user_data:
        user_data = {'username': username, 'email': '', 'phone': ''}
    
    # Get verification status from Firebase
    firebase_ref = db.reference(f'Data/{username}')
    firebase_data = firebase_ref.get() or {}
    
    # Check verification status
    phone_verified = firebase_data.get('phone_verified', False)
    email_verified = firebase_data.get('email_verified', False)
    
    return render_template('email_settings.html', 
                         user=user_data,
                         phone_verified=phone_verified,
                         email_verified=email_verified,
                         **get_user_context())

# --- AI Table Converter Utilities ---
def extract_text_for_table(file_storage):
    filename = file_storage.filename.lower()
    ext = os.path.splitext(filename)[1]
    file_storage.seek(0)
    if ext == '.pdf':
        try:
            reader = PdfReader(file_storage)
            text = "\n".join(page.extract_text() or '' for page in reader.pages)
            return text.strip()
        except Exception as e:
            file_storage.seek(0)
            return f"[PDF extraction error: {e}]"
    elif ext == '.docx':
        try:
            doc = Document(file_storage)
            return "\n".join([p.text for p in doc.paragraphs]).strip()
        except Exception as e:
            file_storage.seek(0)
            return f"[DOCX extraction error: {e}]"
    elif ext == '.txt':
        try:
            file_storage.seek(0)
            return file_storage.read().decode(errors='ignore')
        except Exception as e:
            file_storage.seek(0)
            return f"[TXT extraction error: {e}]"
    else:
        return "[Unsupported file type]"

def convert_unstructured_to_table_with_openrouter(text: str, user_prompt: str = "", target_rows: int = None, target_cols: int = None, retry: bool = False):
    """Call OpenRouter to convert unstructured text to a markdown table with strict constraints.

    Requirements enforced via prompt engineering:
      - Output ONLY a markdown table (pipe syntax) with a single header row.
      - Exact column and row counts as specified.
      - No extra commentary before or after the table.
      - Structured data extraction based on user intent.
    """
    # Build dynamic constraints from requested rows / columns
    constraint_lines = []
    if target_cols:
        constraint_lines.append(f"EXACT number of columns (including header): {target_cols}.")
    else:
        constraint_lines.append("Columns: infer minimal meaningful set (<=8).")
    if target_rows is not None:
        constraint_lines.append(f"EXACT number of DATA rows (excluding header): {target_rows}.")
    else:
        constraint_lines.append("Data rows: concise; avoid duplicates; <= 120.")
    
    retry_note = "CRITICAL: Previous attempt failed validation. You MUST produce exactly the requested structure." if retry else ""
    
    core_instructions = (
        "You are a data extraction expert. Convert the following unstructured data into a markdown table. "
        "Output ONLY the table using pipe (|) syntax - no code fences, no explanations, no additional text.\n"
        + (f"{retry_note}\n" if retry_note else "") +
        "STRICT REQUIREMENTS:\n"
        + "• " + "\n• ".join(constraint_lines) + "\n"
        + "• Header row followed by exactly the requested number of data rows\n"
        + "• Each row must have identical column count\n"
        + "• If data is insufficient, use 'N/A' or relevant placeholder\n"
        + "• Truncate any cell content over 500 characters with '...'\n"
        + "• Focus on extracting the most relevant information\n"
        + "• Do NOT add separator lines (--- | ---) between header and data\n"
        + "• Use standard markdown table format: | Col1 | Col2 |\n"
    )
    
    if user_prompt:
        core_instructions += f"\nUSER SPECIFIC REQUEST: {user_prompt.strip()}\n"
    
    final_prompt = f"{core_instructions}\nSOURCE DATA:\n{text[:15000]}"
    
    headers = {
        'Authorization': f'Bearer {OPENROUTER_API_KEY}',
        'Content-Type': 'application/json',
    }
    data = {
        "model": OPENROUTER_MODEL,
        "messages": [
            {"role": "system", "content": "You are a precise data extraction system that outputs ONLY markdown tables in the exact format requested. Never add explanations, code blocks, or extra text."},
            {"role": "user", "content": final_prompt}
        ],
        "max_tokens": 2000,
        "temperature": 0.1
    }
    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data, timeout=90)
        response.raise_for_status()
        result = response.json()
        content = result.get('choices', [{}])[0].get('message', {}).get('content', '').strip()
        return content or '[No table returned]'
    except Exception as e:
        return f"[OpenRouter API error: {e}]"

def _parse_table_user_directives(user_prompt: str):
    """Extract structural directives (rows/columns) & focus terms from user prompt.

    Returns dict keys:
      target_rows: int|None (# of data rows requested, excludes header)
      target_cols: int|None
      focus_terms: list[str]
    """
    import re
    lower = (user_prompt or '').lower()
    
    # Enhanced patterns to handle text numbers and various phrasings
    number_words = {
        'one': 1, 'two': 2, 'three': 3, 'four': 4, 'five': 5, 'six': 6, 'seven': 7, 'eight': 8, 'nine': 9, 'ten': 10,
        'single': 1, 'double': 2, 'triple': 3
    }
    
    def extract_number(text):
        # Try digit first
        digit_match = re.search(r'\b(\d{1,2})\b', text)
        if digit_match:
            return int(digit_match.group(1))
        # Try word numbers
        for word, num in number_words.items():
            if word in text:
                return num
        return None
    
    # Column patterns - more flexible matching
    col_patterns = [
        r'(?:in|with|exactly|only|make|create|using|format)\s+(?:[\w\s]*?)(\w+)\s*(?:cols?|columns?)',
        r'(\w+)\s*(?:cols?|columns?)(?:\s+(?:and|,))?',
        r'(?:into|as)\s+(?:[\w\s]*?)(\w+)\s*(?:cols?|columns?)'
    ]
    
    row_patterns = [
        r'(?:in|with|exactly|only|make|create|using|format)\s+(?:[\w\s]*?)(\w+)\s*(?:rows?)',
        r'(\w+)\s*(?:rows?)(?:\s+(?:and|,))?',
        r'(?:into|as)\s+(?:[\w\s]*?)(\w+)\s*(?:rows?)'
    ]
    
    target_cols = None
    target_rows = None
    
    # Try column patterns
    for pattern in col_patterns:
        match = re.search(pattern, lower)
        if match:
            target_cols = extract_number(match.group(1))
            if target_cols:
                break
    
    # Try row patterns  
    for pattern in row_patterns:
        match = re.search(pattern, lower)
        if match:
            target_rows = extract_number(match.group(1))
            if target_rows:
                break
    
    # Token extraction for focus terms
    stop = { 'extract','table','format','only','row','rows','column','columns','convert','information','info','data','in','a','the','to','and','as','of','for','make','produce','output','generate','whole','document','all','entire','create','using','into'}
    tokens = [t.strip('.,;:') for t in re.split(r'\s+', lower) if t.strip('.,;:')]
    focus_terms = []
    for t in tokens:
        if len(t) > 3 and t not in stop and t.isalpha() and t not in number_words:
            if t not in focus_terms:
                focus_terms.append(t)
    return { 'target_rows': target_rows, 'target_cols': target_cols, 'focus_terms': focus_terms }

# --- Flask Route: AI Table Converter ---
@app.route('/ai_table_converter', methods=['POST'])
@login_required
def ai_table_converter():
    if 'datafile' not in request.files:
        return jsonify({"error": "No file part in the request."}), 400
    file_storage = request.files['datafile']
    _ = secure_filename(file_storage.filename)
    user_prompt = request.form.get('prompt', '').strip()
    if not user_prompt:
        return jsonify({"error": "Prompt required. Please specify desired rows & columns, e.g. 'Extract in 3 columns and 2 rows'."}), 400
    text = extract_text_for_table(file_storage)
    directives = _parse_table_user_directives(user_prompt)
    target_rows = directives.get('target_rows')
    target_cols = directives.get('target_cols')
    if target_rows is None or target_cols is None:
        return jsonify({"error": "Please include both a row count and a column count in your prompt (e.g. 'in 3 columns and 2 rows')."}), 400
    ai_response = convert_unstructured_to_table_with_openrouter(text, user_prompt, target_rows=target_rows, target_cols=target_cols)
    print('AI Table Converter AI raw response:', ai_response)
    # Parse markdown table
    import re, csv
    table_data = []
    if not ai_response.startswith('['):  # crude error check
        # Extract first markdown table block only
        md_match = re.search(r"((?:\|[^\n]*\|\n?){2,})", ai_response)
        if md_match:
            lines = [ln.strip() for ln in md_match.group(1).splitlines() if ln.strip() and '|' in ln]
            for ln in lines:
                row = [c.strip() for c in ln.strip('|').split('|')]
                table_data.append(row)
            # Remove separator row like --- | ---
            if len(table_data) > 1 and all(re.match(r"^:?-{3,}:?$", c) for c in table_data[1]):
                table_data.pop(1)
        if not table_data:
            # Try CSV fallback
            try:
                for row in csv.reader(ai_response.splitlines()):
                    if row:
                        table_data.append([c.strip() for c in row])
            except Exception:
                pass
    # Sanitize & normalize
    if table_data:
        max_cols = max(len(r) for r in table_data)
        cleaned = []
        for r in table_data:
            # pad or trim
            if len(r) < max_cols:
                r = r + [""] * (max_cols - len(r))
            elif len(r) > max_cols:
                r = r[:max_cols]
            new_row = []
            for cell in r:
                cell = (cell or '').replace('\r', ' ').replace('\n', ' ').strip()
                if len(cell) > 500:
                    cell = cell[:497].rstrip() + '...'
                new_row.append(cell)
            cleaned.append(new_row)
        table_data = cleaned
    else:
        # Represent as a single cell fallback
        table_data = [[ai_response[:500] + ('...' if len(ai_response) > 500 else '')]]
    # Apply user directives (single-row, focus terms)
    # Validate counts; if mismatch do a single retry with stricter flag
    try:
        header = table_data[0] if table_data else []
        data_rows = table_data[1:] if len(table_data) > 1 else []
        mismatch = False
        if len(header) != target_cols:
            mismatch = True
        if len(data_rows) != target_rows:
            mismatch = True
        if mismatch:
            # Retry once with strict flag
            retry_response = convert_unstructured_to_table_with_openrouter(text, user_prompt, target_rows=target_rows, target_cols=target_cols, retry=True)
            if retry_response and not retry_response.startswith('['):
                ai_response2 = retry_response
                # reparse same way
                table_data_retry = []
                md_match2 = re.search(r"((?:\|[^\n]*\|\n?){2,})", ai_response2)
                if md_match2:
                    lines2 = [ln.strip() for ln in md_match2.group(1).splitlines() if ln.strip() and '|' in ln]
                    for ln in lines2:
                        row2 = [c.strip() for c in ln.strip('|').split('|')]
                        table_data_retry.append(row2)
                    if len(table_data_retry) > 1 and all(re.match(r"^:?-{3,}:?$", c) for c in table_data_retry[1]):
                        table_data_retry.pop(1)
                if table_data_retry:
                    table_data = table_data_retry
                    header = table_data[0]
                    data_rows = table_data[1:] if len(table_data) > 1 else []
            # Final validation
            if len(header) != target_cols or len(data_rows) != target_rows:
                return jsonify({
                    "error": f"Model did not comply. Wanted {target_cols} cols & {target_rows} rows, got {len(header)} cols & {len(data_rows)} rows. Please rephrase prompt."})
    except Exception as _val_e:
        print('Validation error:', _val_e)
    # Store in session for later PDF download
    session['table_converter_data'] = table_data
    return jsonify({
        "status": "ok",
        "preview": table_data[:50],  # limit preview rows
        "rows": len(table_data),
        "cols": len(table_data[0]) if table_data else 0
    })

@app.route('/download_table_pdf')
@login_required
def download_table_pdf():
    table_data = session.get('table_converter_data')
    if not table_data:
        return jsonify({"error": "No table data in session. Please convert first."}), 400
    try:
        from fpdf import FPDF
        output = io.BytesIO()
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font('helvetica', 'B', 14)
        pdf.cell(0, 10, 'AI Table Converter Output', ln=True, align='C')
        pdf.ln(6)
        left_margin, right_margin = 15, 15
        effective_page_width = pdf.w - left_margin - right_margin
        pdf.set_left_margin(left_margin)
        pdf.set_right_margin(right_margin)
        # Determine column widths proportional to header length
        header = table_data[0]
        lengths = [max(3, len(h)) for h in header]
        total = sum(lengths)
        col_widths = [effective_page_width * (l / total) for l in lengths]
        line_height = 7
        # Helper to wrap text into lines fitting width
        def wrap_text(txt, width):
            if not txt:
                return ['']
            words = txt.split()
            lines = []
            current = ''
            for w in words:
                test = f"{current} {w}".strip()
                if pdf.get_string_width(test) <= width - 1:
                    current = test
                else:
                    if current:
                        lines.append(current)
                    # If single word too long, hard split
                    if pdf.get_string_width(w) > width - 1:
                        chunk = ''
                        for ch in w:
                            if pdf.get_string_width(chunk + ch) <= width - 1:
                                chunk += ch
                            else:
                                lines.append(chunk)
                                chunk = ch
                        if chunk:
                            current = chunk
                        else:
                            current = ''
                    else:
                        current = w
            if current:
                lines.append(current)
            return lines or ['']
        # Render header
        pdf.set_font('helvetica', 'B', 11)
        max_header_lines = 1
        rows_to_render = table_data
        # Compute height first for header lines
        header_lines = [wrap_text(h, col_widths[i]) for i, h in enumerate(header)]
        max_header_lines = max(len(l) for l in header_lines)
        row_height = line_height * max_header_lines
        y_start = pdf.get_y()
        x_start = pdf.get_x()
        for i, h_lines in enumerate(header_lines):
            x = pdf.get_x()
            y = pdf.get_y()
            pdf.rect(x, y, col_widths[i], row_height)
            for idx, l in enumerate(h_lines):
                pdf.set_xy(x + 1, y + idx * line_height + 1)
                pdf.cell(col_widths[i] - 2, line_height, l, ln=0)
            pdf.set_xy(x + col_widths[i], y)
        pdf.set_xy(x_start, y_start + row_height)
        # Body
        pdf.set_font('helvetica', '', 10)
        for r in rows_to_render[1:]:
            wrapped = [wrap_text(c, col_widths[i]) for i, c in enumerate(r)]
            max_lines = max(len(w) for w in wrapped)
            row_height = line_height * max_lines
            # New page check
            if pdf.get_y() + row_height > pdf.h - 15:
                pdf.add_page()
                pdf.set_font('helvetica', 'B', 11)
                # repeat header on new page
                header_lines = [wrap_text(h, col_widths[i]) for i, h in enumerate(header)]
                max_header_lines = max(len(l) for l in header_lines)
                row_h2 = line_height * max_header_lines
                y_start2 = pdf.get_y()
                x_start2 = pdf.get_x()
                for i, h_lines in enumerate(header_lines):
                    x = pdf.get_x(); y = pdf.get_y()
                    pdf.rect(x, y, col_widths[i], row_h2)
                    for idx, l in enumerate(h_lines):
                        pdf.set_xy(x + 1, y + idx * line_height + 1)
                        pdf.cell(col_widths[i] - 2, line_height, l, ln=0)
                    pdf.set_xy(x + col_widths[i], y)
                pdf.set_xy(x_start2, y_start2 + row_h2)
                pdf.set_font('helvetica', '', 10)
            y_row = pdf.get_y()
            for i, lines in enumerate(wrapped):
                x = pdf.get_x()
                pdf.rect(x, y_row, col_widths[i], row_height)
                for idx, l in enumerate(lines):
                    pdf.set_xy(x + 1, y_row + idx * line_height + 1)
                    pdf.cell(col_widths[i] - 2, line_height, l, ln=0)
                pdf.set_xy(x + col_widths[i], y_row)
            pdf.set_xy(pdf.get_x() - sum(col_widths), y_row + row_height)
        # Generate bytes (force latin1 to avoid Unicode issues for FPDF base fonts)
        pdf_bytes = pdf.output(dest='S').encode('latin1', errors='replace')
        output.write(pdf_bytes)
        output.seek(0)
        out_filename = f"table_conversion_{int(time.time())}.pdf"
        # Upload on-demand now
        try:
            cloudinary_result = cloudinary.uploader.upload(output, resource_type='raw', folder='table_converter', public_id=out_filename)
            url_cloudinary = cloudinary_result.get('secure_url')
            store_url_in_firebase(url_cloudinary, 'table_converter', out_filename)
        except Exception as ue:
            print('Cloudinary upload error (download_table_pdf):', ue)
            return jsonify({"error": f"Cloudinary upload error: {ue}"}), 500
        output.seek(0)
        return send_file(output, as_attachment=True, download_name=out_filename, mimetype='application/pdf')
    except Exception as e:
        print('PDF generation error:', e)
        return jsonify({"error": f"PDF generation error: {e}"}), 500

# --- Download from Cloudinary as attachment ---
@app.route('/download_table/<filename>')
@login_required
def download_table(filename):
    try:
        url = None
        # Get username from session (default to admin)
        username = session.get('username', 'admin')
        safe_key = re.sub(r'[./#$\[\]]', '_', filename)
        ref = db.reference(f'storage/{username}/table_converter/{safe_key}')
        data = ref.get()
        url = data['url'] if data and 'url' in data else None
        if not url:
            url = f"https://res.cloudinary.com/{cloudinary.config().cloud_name}/raw/upload/table_converter/{filename}"
        r = requests.get(url, stream=True)
        r.raise_for_status()
        return send_file(
            io.BytesIO(r.content),
            as_attachment=True,
            download_name=filename,
            mimetype='application/pdf'
        )
    except Exception as e:
        return jsonify({'error': f'Failed to download file: {e}'})

# Render the AI Table Converter page (GET)
@app.route('/ai_table_converter', methods=['GET'])
@login_required
def ai_table_converter_page():
    return render_template('ai_table_converter.html', **get_user_context())

# --- AI PDF Comparison Tool ---

def extract_text_from_pdf_comparison(pdf_file):
    """Extract text from a PDF file object using PyMuPDF."""
    text = []
    try:
        pdf_file.seek(0)
        with fitz.open(stream=pdf_file.read()) as doc:
            for page in doc:
                page_text = page.get_text("text")
                if page_text:
                    text.append(page_text)
    except Exception as e:
        return f"[PDF extraction error: {e}]"
    return "\n\n".join(text).strip()

def call_openrouter_compare(text_a: str, text_b: str, max_chars=20000) -> dict:
    """Call OpenRouter to compare two documents semantically and return structured JSON."""
    # Truncate if very large
    if len(text_a) + len(text_b) > max_chars:
        half = max_chars // 4
        text_a = text_a[:half] + "\n\n...[truncated]...\n\n" + text_a[-half:]
        text_b = text_b[:half] + "\n\n...[truncated]...\n\n" + text_b[-half:]

    system_prompt = (
        "You are an assistant that compares two document versions. "
        "Return a JSON object with two keys: 'summary' (a brief executive summary string) "
        "and 'changes' (an array). Each change must be an object with these fields:\n"
        "  - section: short heading or location (string)\n"
        "  - old: the text in the old document (string) or empty if added\n"
        "  - new: the text in the new document (string) or empty if removed\n"
        "  - status: one of 'added','removed','modified'\n"
        "  - notes: brief natural-language note about the change\n\n"
        "Return valid JSON only, nothing else. Keep notes short (1-2 sentences)."
    )

    user_prompt = (
        "COMPARE DOC A and DOC B.\n\n"
        "DOCUMENT A:\n"
        f"{text_a}\n\n"
        "DOCUMENT B:\n"
        f"{text_b}\n\n"
        "Produce the JSON result as described."
    )

    headers = {
        'Authorization': f'Bearer {OPENROUTER_API_KEY}',
        'Content-Type': 'application/json',
    }
    data = {
        "model": OPENROUTER_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "max_tokens": 3000,
        "temperature": 0.0
    }

    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data, timeout=120)
        response.raise_for_status()
        result = response.json()
        content = result.get('choices', [{}])[0].get('message', {}).get('content', '').strip()
        
        # Extract JSON from response
        start = content.find("{")
        end = content.rfind("}")
        if start != -1 and end != -1:
            json_text = content[start:end+1]
        else:
            json_text = content
        parsed = json.loads(json_text)
        return parsed
    except Exception as e:
        # Fallback to simple diff
        return fallback_text_diff(text_a, text_b)

def fallback_text_diff(text_a: str, text_b: str, max_changes=50) -> dict:
    """Simple fallback: use difflib to generate line-based differences."""
    a_lines = text_a.splitlines()
    b_lines = text_b.splitlines()
    diff = difflib.unified_diff(a_lines, b_lines, lineterm="")
    changes = []
    added_count = removed_count = 0
    
    for i, line in enumerate(diff):
        if line.startswith("---") or line.startswith("+++") or line.startswith("@@"):
            continue
        if line.startswith("- "):
            removed_count += 1
            changes.append({
                "section": f"Line {i}",
                "old": line[2:],
                "new": "",
                "status": "removed",
                "notes": "Line removed"
            })
        elif line.startswith("+ "):
            added_count += 1
            changes.append({
                "section": f"Line {i}",
                "old": "",
                "new": line[2:],
                "status": "added",
                "notes": "Line added"
            })
        if len(changes) >= max_changes:
            break

    summary = f"Found {added_count} additions and {removed_count} deletions between documents."
    return {"summary": summary, "changes": changes}

def generate_comparison_pdf(metadata: dict, comparisons: list) -> io.BytesIO:
    """Generate a PDF comparison report using reportlab with proper text wrapping."""
    output = io.BytesIO()
    doc = SimpleDocTemplate(output, pagesize=A4, rightMargin=30, leftMargin=30, topMargin=30, bottomMargin=30)
    styles = getSampleStyleSheet()
    elements = []

    # Title
    title = Paragraph(metadata.get("title", "DocShift AI Comparison Report"), styles['Title'])
    elements.append(title)
    elements.append(Spacer(1, 12))

    # Metadata
    meta_text = f"Old file: {metadata.get('file_a_name','')} &nbsp;&nbsp;&nbsp; New file: {metadata.get('file_b_name','')}<br/>Generated: {metadata.get('generated_at','')}"
    elements.append(Paragraph(meta_text, styles['Normal']))
    elements.append(Spacer(1, 12))

    # Executive summary
    elements.append(Paragraph("<b>Executive Summary</b>", styles['Heading2']))
    elements.append(Paragraph(metadata.get("summary", ""), styles['Normal']))
    elements.append(Spacer(1, 12))

    elements.append(Paragraph("<b>Detailed Changes</b>", styles['Heading2']))
    elements.append(Spacer(1, 8))

    # Custom style for table cells with smaller font and proper wrapping
    cell_style = ParagraphStyle(
        'CellText',
        parent=styles['Normal'],
        fontSize=8,
        leading=10,
        wordWrap='LTR',
        alignment=0,  # Left align
        spaceAfter=3,
    )

    # Helper function to create wrapped paragraphs for table cells
    def create_cell_paragraph(text, max_chars=200):
        if not text:
            return Paragraph("", cell_style)
        
        # Truncate very long text but preserve word boundaries
        if len(text) > max_chars:
            words = text.split()
            truncated_words = []
            char_count = 0
            
            for word in words:
                if char_count + len(word) + 1 <= max_chars - 3:  # Leave space for "..."
                    truncated_words.append(word)
                    char_count += len(word) + 1
                else:
                    break
            
            text = " ".join(truncated_words) + "..."
        
        # Escape HTML characters and create paragraph
        escaped_text = html.escape(text)
        return Paragraph(escaped_text, cell_style)

    # Table header with Paragraph objects
    header_style = ParagraphStyle(
        'HeaderText',
        parent=styles['Normal'],
        fontSize=10,
        leading=12,
        wordWrap='LTR',
        alignment=1,  # Center align
        textColor=colors.black,
        fontName='Helvetica-Bold'
    )

    table_data = [[
        Paragraph("Section", header_style),
        Paragraph("Status", header_style), 
        Paragraph("Old Content", header_style),
        Paragraph("New Content", header_style)
    ]]
    
    # Add comparison rows with properly wrapped content
    for ch in comparisons:
        # Create section cell with truncation
        section_text = ch.get("section", "")
        if len(section_text) > 40:
            section_text = section_text[:37] + "..."
        section_para = Paragraph(section_text, cell_style)
        
        # Create status cell with color coding
        status = ch.get("status", "").upper()
        status_style = ParagraphStyle(
            'StatusText',
            parent=cell_style,
            fontSize=8,
            fontName='Helvetica-Bold',
            alignment=1,  # Center align
        )
        
        # Color code the status
        if status == "ADDED":
            status_style.textColor = colors.green
            status_text = "✓ ADDED"
        elif status == "REMOVED": 
            status_style.textColor = colors.red
            status_text = "✗ REMOVED"
        elif status == "MODIFIED":
            status_style.textColor = colors.orange
            status_text = "⚠ MODIFIED"
        else:
            status_text = status
            
        status_para = Paragraph(status_text, status_style)
        
        # Create content cells with proper wrapping
        old_para = create_cell_paragraph(ch.get("old", ""), 180)
        new_para = create_cell_paragraph(ch.get("new", ""), 180)
        
        table_data.append([section_para, status_para, old_para, new_para])

    # Create table with proper column widths
    table = Table(table_data, colWidths=[90, 70, 190, 190], repeatRows=1)
    style = TableStyle([
        # Header styling
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#e8e8e8")),
        ('TEXTCOLOR', (0,0), (-1,0), colors.black),
        ('ALIGN', (0,0), (-1,0), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        
        # Cell styling
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,0), 10),
        ('BOTTOMPADDING', (0,0), (-1,0), 8),
        ('TOPPADDING', (0,0), (-1,0), 8),
        
        # Grid and borders
        ('GRID', (0,0), (-1,-1), 0.5, colors.grey),
        ('LINEBELOW', (0,0), (-1,0), 1, colors.black),
        
        # Row backgrounds
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor("#f9f9f9")]),
        
        # Cell padding
        ('LEFTPADDING', (0,0), (-1,-1), 6),
        ('RIGHTPADDING', (0,0), (-1,-1), 6),
        ('TOPPADDING', (0,1), (-1,-1), 8),
        ('BOTTOMPADDING', (0,1), (-1,-1), 8),
    ])
    table.setStyle(style)
    elements.append(table)
    elements.append(Spacer(1, 12))

    # Summary footer
    if comparisons:
        total_changes = len(comparisons)
        elements.append(Paragraph(f"<b>Total Changes Found:</b> {total_changes}", styles['Normal']))
        elements.append(Spacer(1, 8))

    doc.build(elements)
    output.seek(0)
    return output

def generate_comparison_docx(metadata: dict, comparisons: list) -> io.BytesIO:
    """Generate a DOCX comparison report."""
    from docx import Document
    output = io.BytesIO()
    doc = Document()
    
    # Title
    title = doc.add_heading(metadata.get("title", "DocShift AI Comparison Report"), level=1)
    title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

    # Metadata
    meta_para = doc.add_paragraph()
    meta_para.add_run(f"Old file: {metadata.get('file_a_name', '')}\n")
    meta_para.add_run(f"New file: {metadata.get('file_b_name', '')}\n")
    meta_para.add_run(f"Generated: {metadata.get('generated_at', '')}\n")

    doc.add_paragraph("")

    # Executive Summary
    doc.add_heading("Executive Summary", level=2)
    doc.add_paragraph(metadata.get("summary", ""))

    doc.add_paragraph("")
    doc.add_heading("Detailed Changes", level=2)

    # Changes
    for change in comparisons:
        tbl = doc.add_table(rows=2, cols=3)
        hdr_cells = tbl.rows[0].cells
        hdr_cells[0].text = "Section"
        hdr_cells[1].text = "Old"
        hdr_cells[2].text = "New"

        row_cells = tbl.rows[1].cells
        row_cells[0].text = change.get("section", "")
        row_cells[1].text = change.get("old", "")
        row_cells[2].text = change.get("new", "")

        p = doc.add_paragraph()
        p.add_run("Status: ").bold = True
        p.add_run(change.get("status", ""))
        if change.get("notes"):
            p.add_run("\nNotes: ").bold = True
            p.add_run(change.get("notes", ""))
        doc.add_paragraph("")

    doc.save(output)
    output.seek(0)
    return output

@app.route('/ai_pdf_comparison', methods=['GET'])
@login_required
def ai_pdf_comparison_page():
    return render_template('ai_pdf_comparison.html', **get_user_context())

@app.route('/ai_pdf_comparison', methods=['POST'])
@login_required
def ai_pdf_comparison():
    if 'file_a' not in request.files or 'file_b' not in request.files:
        return jsonify({"error": "Both PDF files are required."}), 400
    
    file_a = request.files['file_a']
    file_b = request.files['file_b']
    
    if not file_a.filename or not file_b.filename:
        return jsonify({"error": "Please select both PDF files."}), 400

    try:
        # Extract text from both PDFs
        text_a = extract_text_from_pdf_comparison(file_a)
        text_b = extract_text_from_pdf_comparison(file_b)
        
        if text_a.startswith('[PDF extraction error') or text_b.startswith('[PDF extraction error'):
            return jsonify({"error": "Failed to extract text from one or both PDF files."}), 400

        # Compare using OpenRouter
        comparison_result = call_openrouter_compare(text_a, text_b)
        summary = comparison_result.get("summary", "")
        changes = comparison_result.get("changes", [])

        # Generate reports
        timestamp = datetime.utcnow().strftime("%Y-%m-%d_%H%M%S")
        metadata = {
            "title": "DocShift AI Comparison Report",
            "file_a_name": secure_filename(file_a.filename),
            "file_b_name": secure_filename(file_b.filename),
            "generated_at": timestamp,
            "summary": summary
        }

        # Store in session for download
        session['pdf_comparison_data'] = {
            'metadata': metadata,
            'changes': changes
        }

        return jsonify({
            "status": "ok",
            "summary": summary,
            "change_count": len(changes),
            "preview_changes": changes[:20],  # Preview first 20 changes
            "total_changes": len(changes),
            "metadata": {
                'file_a_name': metadata.get('file_a_name', ''),
                'file_b_name': metadata.get('file_b_name', ''),
                'generated_at': metadata.get('generated_at', '')
            }
        })

    except Exception as e:
        print('PDF comparison error:', e)
        return jsonify({"error": f"Comparison failed: {str(e)}"}), 500

@app.route('/download_comparison_pdf')
@login_required
def download_comparison_pdf():
    comparison_data = session.get('pdf_comparison_data')
    if not comparison_data:
        return jsonify({"error": "No comparison data found. Please run comparison first."}), 400

    try:
        metadata = comparison_data['metadata']
        changes = comparison_data['changes']
        
        # Generate PDF with error handling
        try:
            pdf_output = generate_comparison_pdf(metadata, changes)
        except Exception as pdf_error:
            print('PDF generation specific error:', pdf_error)
            import traceback
            traceback.print_exc()
            return jsonify({"error": f"PDF generation failed: {str(pdf_error)}"}), 500
        
        out_filename = f"comparison_report.pdf"
        
        # Upload to Cloudinary
        try:
            pdf_output.seek(0)  # Reset buffer position
            cloudinary_result = cloudinary.uploader.upload(
                pdf_output, 
                resource_type='raw', 
                folder='pdf_comparison', 
                public_id=out_filename
            )
            url_cloudinary = cloudinary_result.get('secure_url')
            store_url_in_firebase(url_cloudinary, 'pdf_comparison', out_filename)
        except Exception as ue:
            print('Cloudinary upload error (comparison PDF):', ue)
            # Continue with direct download even if upload fails
        
        pdf_output.seek(0)  # Reset buffer position for download
        return send_file(
            pdf_output, 
            as_attachment=True, 
            download_name=out_filename, 
            mimetype='application/pdf'
        )
    except Exception as e:
        print('PDF generation error:', e)
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"PDF generation error: {str(e)}"}), 500

@app.route('/download_comparison_docx')
@login_required
def download_comparison_docx():
    comparison_data = session.get('pdf_comparison_data')
    if not comparison_data:
        return jsonify({"error": "No comparison data found. Please run comparison first."}), 400

    try:
        metadata = comparison_data['metadata']
        changes = comparison_data['changes']
        
        # Generate DOCX
        docx_output = generate_comparison_docx(metadata, changes)
        out_filename = f"comparison_{metadata['generated_at']}.docx"
        
        # Upload to Cloudinary
        try:
            cloudinary_result = cloudinary.uploader.upload(
                docx_output, 
                resource_type='raw', 
                folder='pdf_comparison', 
                public_id=out_filename
            )
            url_cloudinary = cloudinary_result.get('secure_url')
            store_url_in_firebase(url_cloudinary, 'pdf_comparison', out_filename)
        except Exception as ue:
            print('Cloudinary upload error (comparison DOCX):', ue)
            return jsonify({"error": f"Upload error: {ue}"}), 500
        
        docx_output.seek(0)
        return send_file(
            docx_output, 
            as_attachment=True, 
            download_name=out_filename, 
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
    except Exception as e:
        print('DOCX generation error:', e)
        return jsonify({"error": f"DOCX generation error: {e}"}), 500


# ==================== ADMIN DASHBOARD ENDPOINTS ====================
init_db()  # Initialize DB once on startup
update_users_table()  # Update users table schema
# --- Run Flask App ---

if __name__ == '__main__':
    
    print("🚀 Starting DocShift on 127.0.0.1:5000")
    print("� Email verification enabled")
    print("📱 Phone/SMS verification DISABLED")
    print("🔥 Firebase SMS functionality has been removed")
    port = int(os.environ.get('PORT', 8080))
    app.run(debug=False, host='0.0.0.0', port=port)